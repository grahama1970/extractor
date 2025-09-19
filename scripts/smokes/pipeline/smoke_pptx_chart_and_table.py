#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.10"
# dependencies = [
#   "typer>=0.12",
#   "python-pptx>=0.6.22",
# ]
# ///
"""PPTX: detect both a chart and a table shape in Stage 07 (as Figure and Table)."""

from __future__ import annotations

import json
from pathlib import Path
import typer
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from extractor.pipeline.structured_pipeline import STRUCTURED_PIPELINES, run_structured_pipeline
from extractor.core.providers.pptx import PPTXProvider

app = typer.Typer(add_completion=False)


@app.command()
def main(tmp_dir: Path = typer.Option(Path("data/results/structured_parity_smoke/pptx_synth"))):
    tmp_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = tmp_dir / "chart_table.pptx"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Add table
    x, y, cx, cy = Inches(1), Inches(1), Inches(3), Inches(1)
    table_shape = slide.shapes.add_table(2, 2, x, y, cx, cy)
    table = table_shape.table
    table.cell(0, 0).text = "A"; table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"; table.cell(1, 1).text = "2"
    # Add chart
    chart_data = ChartData(); chart_data.categories = ['Q1', 'Q2']; chart_data.add_series('S1', (1, 2))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(3), Inches(4), Inches(3), chart_data)
    prs.save(pptx_path)

    meta = STRUCTURED_PIPELINES[PPTXProvider]
    artifacts = run_structured_pipeline(PPTXProvider, pptx_path, tmp_dir, stage_prefix=meta.stage_prefix, skip_export10=True, skip_embeddings10=True, fast_embeddings10=True)
    s07 = json.loads(Path(artifacts["stage07"]).read_text())
    sections = s07.get("reflowed_sections") or []
    if not sections:
        typer.echo("No sections built in PPTX chart+table synthetic.", err=True)
        raise typer.Exit(code=1)
    # Tolerate either figure or image for chart; require at least one figure and one table across sections
    has_fig = any((len(s.get("figures") or []) > 0) for s in sections)
    has_tab = any((len(s.get("tables") or []) > 0) for s in sections)
    if not has_tab:
        typer.echo("No tables detected in PPTX Stage 07.", err=True)
        raise typer.Exit(code=1)
    if not has_fig:
        typer.echo("No figures detected in PPTX Stage 07 (chart).", err=True)
        raise typer.Exit(code=1)
    typer.echo("PPTX chart+table detection passed.")


if __name__ == "__main__":
    app()

