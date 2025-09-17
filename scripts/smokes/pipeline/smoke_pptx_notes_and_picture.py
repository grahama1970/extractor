#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.10"
# dependencies = [
#   "typer>=0.12",
#   "python-pptx>=0.6.22",
# ]
# ///
"""PPTX: notes_slide content maps into reflowed_text; picture yields a figure in Stage 07.

Creates a synthetic PPTX with one slide, a note, and a tiny PNG picture.
"""

from __future__ import annotations

import base64
import json
from pathlib import Path
import typer
from pptx import Presentation
from pptx.util import Inches

from extractor.pipeline.structured_pipeline import STRUCTURED_PIPELINES, run_structured_pipeline
from extractor.core.providers.pptx import PPTXProvider

app = typer.Typer(add_completion=False)


@app.command()
def main(tmp_dir: Path = typer.Option(Path("data/results/structured_parity_smoke/pptx_synth"))):
    tmp_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = tmp_dir / "notes_pic.pptx"
    png_path = tmp_dir / "tiny.png"
    png_b64 = b"iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR4nGNgYAAAAAMAASsJTYQAAAAASUVORK5CYII="
    png_path.write_bytes(base64.b64decode(png_b64))

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "T" if slide.shapes.title else None
    # Add notes
    notes = slide.notes_slide
    notes.notes_text_frame.text = "Speaker note here"
    # Add picture
    slide.shapes.add_picture(str(png_path), Inches(1), Inches(1), width=Inches(1), height=Inches(1))
    prs.save(pptx_path)

    meta = STRUCTURED_PIPELINES[PPTXProvider]
    artifacts = run_structured_pipeline(PPTXProvider, pptx_path, tmp_dir, stage_prefix=meta.stage_prefix, skip_export10=True, skip_embeddings10=True, fast_embeddings10=True)
    s07 = json.loads(Path(artifacts["stage07"]).read_text())
    secs = s07.get("reflowed_sections") or []
    assert secs, "No sections built"
    text = secs[0].get("reflowed_text") or ""
    if "Speaker note here" not in text:
        typer.echo("PPTX notes not present in reflowed_text.", err=True)
        raise typer.Exit(code=1)
    figs = secs[0].get("figures") or []
    if not figs:
        typer.echo("PPTX picture not captured as figure/image in Stage 07.", err=True)
        raise typer.Exit(code=1)
    typer.echo("PPTX notes + picture mapping passed.")


if __name__ == "__main__":
    app()

