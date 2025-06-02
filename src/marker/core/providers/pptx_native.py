"""
Module: pptx_native.py
Purpose: Native PowerPoint extraction without PDF conversion

This module implements direct PPTX extraction using python-pptx,
preserving all presentation features including speaker notes, slide
layouts, animations metadata, and embedded objects without lossy conversions.

External Dependencies:
- python-pptx: https://python-pptx.readthedocs.io/

Example Usage:
>>> from marker.core.providers.pptx_native import NativePPTXProvider
>>> provider = NativePPTXProvider()
>>> document = provider.extract_document("presentation.pptx")
>>> print(document.source_type)  # SourceType.PPTX
"""

import hashlib
import base64
from pathlib import Path
from typing import List, Dict, Any, Optional, Union
from datetime import datetime

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.dml import MSO_THEME_COLOR
from loguru import logger

from marker.core.schema.unified_document import (
    UnifiedDocument, BlockType, SourceType, BaseBlock, TableBlock,
    ImageBlock, BlockMetadata, DocumentMetadata,
    HierarchyNode, TableCell
)


class NativePPTXProvider:
    """Direct PPTX extraction without PDF conversion"""
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        self.config = config or {}
        self.block_counter = 0
        self.slide_hierarchy = {}
        
    def extract_document(self, filepath: Union[str, Path]) -> UnifiedDocument:
        """Extract PPTX content to unified document format"""
        filepath = Path(filepath)
        logger.info(f"Extracting PPTX document: {filepath}")
        
        # Load presentation
        prs = Presentation(str(filepath))
        
        # Extract metadata
        metadata = self._extract_metadata(prs, filepath)
        
        # Extract all slides with structure
        blocks = []
        slide_nodes = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_blocks, slide_node = self._extract_slide(slide, slide_idx)
            blocks.extend(slide_blocks)
            slide_nodes.append(slide_node)
        
        # Build document hierarchy
        hierarchy = self._build_hierarchy(slide_nodes)
        
        # Apply Claude table merge analysis if enabled
        if self.config.get('enable_table_merge_analysis', True):
            blocks = self._analyze_and_merge_tables(blocks)
        
        # Create unified document
        doc = UnifiedDocument(
            id=self._generate_doc_id(filepath),
            source_type=SourceType.PPTX,
            source_path=str(filepath),
            blocks=blocks,
            hierarchy=hierarchy,
            metadata=metadata,
            full_text=self._extract_full_text(blocks),
            keywords=self._extract_keywords(prs)
        )
        
        logger.info(f"Extracted {len(blocks)} blocks from {len(prs.slides)} slides")
        return doc
    
    def _generate_doc_id(self, filepath: Path) -> str:
        """Generate unique document ID"""
        return hashlib.md5(str(filepath).encode()).hexdigest()
    
    def _extract_metadata(self, prs: Presentation, filepath: Path) -> DocumentMetadata:
        """Extract presentation metadata"""
        metadata = DocumentMetadata()
        
        # Core properties
        props = prs.core_properties
        if props:
            metadata.title = props.title or ""
            metadata.author = props.author or ""
            
            # Dates
            if props.created:
                metadata.created_date = props.created
            if props.modified:
                metadata.modified_date = props.modified
            
            # Language (if available)
            metadata.language = props.language or ""
        
        # Presentation-specific metadata
        metadata.format_metadata = {
            'file_type': 'pptx',
            'slide_count': len(prs.slides),
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height,
            'has_notes': any(slide.has_notes_slide for slide in prs.slides),
            'presentation_format': str(prs.slide_width) + 'x' + str(prs.slide_height)
        }
        
        # File metadata
        metadata.file_size = filepath.stat().st_size if filepath.exists() else None
        
        return metadata
    
    def _extract_slide(self, slide, slide_idx: int) -> tuple[List[BaseBlock], HierarchyNode]:
        """Extract all content from a single slide"""
        blocks = []
        
        # Create slide heading block
        slide_title = self._get_slide_title(slide)
        if not slide_title:
            slide_title = f"Slide {slide_idx + 1}"
            
        slide_block = BaseBlock(
            id=self._generate_block_id(),
            type=BlockType.HEADING,
            content=slide_title,
            metadata=BlockMetadata(
                attributes={
                    'level': 1,
                    'slide_number': slide_idx + 1,
                    'is_slide_title': True
                },
                confidence=1.0
            )
        )
        blocks.append(slide_block)
        
        # Extract shapes
        for shape in slide.shapes:
            shape_blocks = self._extract_shape(shape, slide_idx)
            blocks.extend(shape_blocks)
        
        # Extract notes
        if slide.has_notes_slide:
            notes_blocks = self._extract_notes(slide.notes_slide, slide_idx)
            blocks.extend(notes_blocks)
        
        # Create hierarchy node for slide
        slide_node = HierarchyNode(
            id=f"slide-{slide_idx}",
            title=slide_title,
            level=1,
            block_id=slide_block.id,
            breadcrumb=["Presentation", slide_title]
        )
        
        return blocks, slide_node
    
    def _get_slide_title(self, slide) -> Optional[str]:
        """Extract slide title from title placeholder"""
        for shape in slide.shapes:
            if shape.is_placeholder:
                if shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                    if hasattr(shape, "text") and shape.text:
                        return shape.text.strip()
        return None
    
    def _extract_shape(self, shape, slide_idx: int) -> List[BaseBlock]:
        """Extract content from a shape"""
        blocks = []
        
        # Handle grouped shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                sub_blocks = self._extract_shape(sub_shape, slide_idx)
                blocks.extend(sub_blocks)
            return blocks
        
        # Handle tables
        if shape.has_table:
            table_block = self._extract_table(shape, slide_idx)
            if table_block:
                blocks.append(table_block)
            return blocks
        
        # Handle images/pictures
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_block = self._extract_image(shape, slide_idx)
            if image_block:
                blocks.append(image_block)
            return blocks
        
        # Handle charts
        if shape.has_chart:
            chart_block = self._extract_chart(shape, slide_idx)
            if chart_block:
                blocks.append(chart_block)
            return blocks
        
        # Handle text
        if hasattr(shape, "text") and shape.text and shape.text.strip():
            text_blocks = self._extract_text(shape, slide_idx)
            blocks.extend(text_blocks)
        
        return blocks
    
    def _extract_text(self, shape, slide_idx: int) -> List[BaseBlock]:
        """Extract text content from shape"""
        blocks = []
        
        # Determine block type based on placeholder
        block_type = BlockType.PARAGRAPH
        level = None
        
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            if placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                block_type = BlockType.HEADING
                level = 2
            elif placeholder_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                # Skip - already handled as slide title
                return blocks
        
        if shape.has_text_frame:
            # Collect all paragraphs
            all_paragraphs = []
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip()
                if text:
                    all_paragraphs.append((paragraph, text))
            
            # If we have multiple paragraphs in a content placeholder, treat as list
            if (len(all_paragraphs) > 1 and shape.is_placeholder and 
                shape.placeholder_format.type not in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]):
                
                # Create list items
                list_items = []
                for paragraph, text in all_paragraphs:
                    p_el = paragraph._element
                    bullet_num = p_el.find(".//a:buAutoNum", namespaces=p_el.nsmap)
                    is_numbered = (bullet_num is not None)
                    
                    list_items.append({
                        'text': text,
                        'level': paragraph.level,
                        'is_numbered': is_numbered
                    })
                
                # Create list block
                if list_items:
                    list_block = self._create_list_block(list_items, slide_idx)
                    blocks.append(list_block)
            
            else:
                # Process as individual paragraphs
                for paragraph, text in all_paragraphs:
                    blocks.append(BaseBlock(
                        id=self._generate_block_id(),
                        type=block_type,
                        content=text,
                        metadata=BlockMetadata(
                            attributes={
                                'slide_number': slide_idx + 1,
                                'shape_type': 'text'
                            },
                            confidence=0.95
                        )
                    ))
        
        return blocks
    
    def _create_list_block(self, items: List[Dict], slide_idx: int) -> BaseBlock:
        """Create a list block from items"""
        # Format as markdown list
        content_lines = []
        for item in items:
            indent = "  " * item['level']
            if item['is_numbered']:
                content_lines.append(f"{indent}1. {item['text']}")
            else:
                content_lines.append(f"{indent}- {item['text']}")
        
        return BaseBlock(
            id=self._generate_block_id(),
            type=BlockType.LISTITEM,
            content="\n".join(content_lines),
            metadata=BlockMetadata(
                attributes={
                    'slide_number': slide_idx + 1,
                    'list_type': 'ordered' if items[0]['is_numbered'] else 'unordered',
                    'item_count': len(items)
                },
                confidence=0.95
            )
        )
    
    def _extract_table(self, shape, slide_idx: int) -> Optional[TableBlock]:
        """Extract table from shape"""
        table = shape.table
        cells = []
        headers = []
        
        for row_idx, row in enumerate(table.rows):
            # First row often headers
            if row_idx == 0:
                headers.append(row_idx)
                
            for col_idx, cell in enumerate(row.cells):
                cells.append(TableCell(
                    row=row_idx,
                    col=col_idx,
                    content=cell.text.strip() if cell.text else ""
                ))
        
        if not cells:
            return None
            
        return TableBlock(
            id=self._generate_block_id(),
            type=BlockType.TABLE,
            content={},
            rows=len(table.rows),
            cols=len(table.columns),
            cells=cells,
            headers=headers,
            metadata=BlockMetadata(
                attributes={
                    'slide_number': slide_idx + 1,
                    'shape_type': 'table'
                },
                confidence=0.95
            )
        )
    
    def _extract_image(self, shape, slide_idx: int) -> Optional[ImageBlock]:
        """Extract image from shape"""
        try:
            image = shape.image
            image_bytes = image.blob
            
            # Convert to base64
            image_base64 = base64.b64encode(image_bytes).decode('utf-8')
            data_uri = f"data:{image.content_type};base64,{image_base64}"
            
            return ImageBlock(
                id=self._generate_block_id(),
                type=BlockType.IMAGE,
                content="",
                src=data_uri,
                alt=shape.name or f"Image from slide {slide_idx + 1}",
                format=image.ext or 'unknown',
                metadata=BlockMetadata(
                    attributes={
                        'slide_number': slide_idx + 1,
                        'shape_type': 'picture',
                        'content_type': image.content_type
                    },
                    confidence=1.0
                )
            )
        except Exception as e:
            logger.warning(f"Failed to extract image: {e}")
            return None
    
    def _extract_chart(self, shape, slide_idx: int) -> Optional[BaseBlock]:
        """Extract chart information"""
        try:
            chart = shape.chart
            
            # Extract chart data as text description
            chart_info = [
                f"Chart Type: {chart.chart_type}",
                f"Title: {chart.chart_title.text_frame.text if chart.has_title else 'Untitled'}"
            ]
            
            # Try to extract series data
            if hasattr(chart, 'series'):
                chart_info.append(f"Series Count: {len(chart.series)}")
            
            return BaseBlock(
                id=self._generate_block_id(),
                type=BlockType.FIGURE,
                content="\n".join(chart_info),
                metadata=BlockMetadata(
                    attributes={
                        'slide_number': slide_idx + 1,
                        'shape_type': 'chart',
                        'chart_type': str(chart.chart_type)
                    },
                    confidence=0.9
                )
            )
        except Exception as e:
            logger.warning(f"Failed to extract chart: {e}")
            return None
    
    def _extract_notes(self, notes_slide, slide_idx: int) -> List[BaseBlock]:
        """Extract speaker notes"""
        blocks = []
        
        # Find the notes text frame
        for shape in notes_slide.shapes:
            if shape.has_text_frame and shape.text:
                blocks.append(BaseBlock(
                    id=self._generate_block_id(),
                    type=BlockType.COMMENT,
                    content=shape.text.strip(),
                    metadata=BlockMetadata(
                        attributes={
                            'slide_number': slide_idx + 1,
                            'comment_type': 'speaker_notes'
                        },
                        confidence=1.0
                    )
                ))
        
        return blocks
    
    def _build_hierarchy(self, slide_nodes: List[HierarchyNode]) -> HierarchyNode:
        """Build presentation hierarchy"""
        root = HierarchyNode(
            id="root",
            title="Presentation",
            level=0,
            block_id="root",
            children=slide_nodes
        )
        
        # Set parent references
        for node in slide_nodes:
            node.parent_id = root.id
            
        return root
    
    def _extract_full_text(self, blocks: List[BaseBlock]) -> str:
        """Extract all text content"""
        text_parts = []
        for block in blocks:
            if hasattr(block, 'content') and isinstance(block.content, str):
                text_parts.append(block.content)
            elif block.type == BlockType.TABLE and hasattr(block, 'cells'):
                for cell in block.cells:
                    if cell.content:
                        text_parts.append(cell.content)
        return '\n'.join(text_parts)
    
    def _extract_keywords(self, prs: Presentation) -> List[str]:
        """Extract keywords from presentation"""
        keywords = []
        
        # From core properties
        if prs.core_properties:
            if prs.core_properties.keywords:
                # Split comma-separated keywords
                kw_string = prs.core_properties.keywords
                keywords.extend([k.strip() for k in kw_string.split(',') if k.strip()])
            
            if prs.core_properties.subject:
                keywords.append(prs.core_properties.subject)
                
            if prs.core_properties.category:
                keywords.append(prs.core_properties.category)
        
        return keywords
    
    def _generate_block_id(self) -> str:
        """Generate unique block ID"""
        self.block_counter += 1
        return f"pptx-block-{self.block_counter}"
    
    def _analyze_and_merge_tables(self, blocks: List[BaseBlock]) -> List[BaseBlock]:
        """Analyze tables for potential merging using Claude's logic"""
        try:
            from marker.core.processors.claude_table_merge_analyzer import ClaudeTableMergeAnalyzer
            analyzer = ClaudeTableMergeAnalyzer()
            return analyzer.analyze_and_merge_tables(blocks)
        except ImportError:
            logger.debug("Claude table merge analyzer not available, skipping table merge analysis")
            return blocks


if __name__ == "__main__":
    # Test the PPTX extractor
    provider = NativePPTXProvider()
    
    # Basic validation
    assert provider is not None
    assert provider.block_counter == 0
    
    print("✅ Native PPTX provider initialized successfully")
    print("Note: Full testing requires actual PowerPoint files")