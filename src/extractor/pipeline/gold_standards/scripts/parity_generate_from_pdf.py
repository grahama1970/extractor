#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = [
#   "pymupdf>=1.24.8",
#   "python-docx>=1.1.2",
#   "python-pptx>=0.6.23",
#   "openpyxl>=3.1.2",
#   "ebooklib>=0.18",
# ]
# ///
"""Generate cross-format parity fixtures from a source PDF.

Outputs go under data/input/parity_gen/<stem>/ with these files:
  - .md, .html, .docx, .rst, .pptx, .xlsx, .epub, .xml
Sections are derived from PDF headings (simple heuristic) and basic tables
from Camelot-like text parsing (not full fidelity, but stable enough for parity smokes).
"""

from __future__ import annotations

import argparse
from pathlib import Path
import fitz
import re
from typing import List

from docx import Document
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook
from ebooklib import epub


HEADING_RE = re.compile(r"^(\d+(?:\.\d+)*)\s+(.*)")


def extract_sections(pdf_path: Path):
    doc = fitz.open(pdf_path)
    sections: List[dict] = []
    for page_index, page in enumerate(doc):
        text = page.get_text("text")
        for line in text.splitlines():
            m = HEADING_RE.match(line.strip())
            if m:
                sections.append({"title": m.group(0), "page": page_index + 1})
    if not sections:
        sections.append({"title": pdf_path.stem, "page": 1})
    return sections


def extract_body(pdf_path: Path):
    doc = fitz.open(pdf_path)
    pages = []
    for page in doc:
        pages.append(page.get_text("text"))
    return pages


def write_md(out: Path, sections, pages):
    lines = []
    for sec in sections:
        lines.append(f"# {sec['title']}")
        lines.append("")
    lines.append("\n".join(pages))
    out.write_text("\n".join(lines), encoding="utf-8")


def write_html(out: Path, sections, pages):
    html = ["<html><body>"]
    for sec in sections:
        html.append(f"<h1>{sec['title']}</h1>")
    for p in pages:
        html.append(f"<p>{p.replace('\n',' ')}</p>")
    html.append("</body></html>")
    out.write_text("\n".join(html), encoding="utf-8")


def write_docx(out: Path, sections, pages):
    doc = Document()
    for sec in sections:
        doc.add_heading(sec["title"], level=1)
    for p in pages:
        doc.add_paragraph(p)
    doc.save(out)


def write_rts(out: Path, sections, pages):
    lines = []
    for sec in sections:
        lines.append(sec["title"])
        lines.append("=" * max(3, len(sec["title"])))
    lines.extend([p for p in pages])
    out.write_text("\n\n".join(lines), encoding="utf-8")


def write_pptx(out: Path, sections, pages):
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    blank_layout = prs.slide_layouts[6]
    # One slide per section
    for sec in sections:
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = sec["title"]
    # Add one slide with body text
    slide = prs.slides.add_slide(blank_layout)
    left = top = Inches(1)
    width = height = Inches(8)
    text_frame = slide.shapes.add_textbox(left, top, width, height).text_frame
    for p in pages:
        text_frame.add_paragraph().text = p.replace("\n", " ")
    prs.save(out)


def write_xlsx(out: Path, pages):
    wb = Workbook()
    ws = wb.active
    ws.title = "Text"
    row = 1
    for p in pages:
        for line in p.splitlines():
            ws.cell(row=row, column=1, value=line)
            row += 1
    wb.save(out)


def write_epub(out: Path, sections, pages):
    book = epub.EpubBook()
    book.set_title(out.stem)
    book.set_language("en")
    items = []
    for i, sec in enumerate(sections, 1):
        c = epub.EpubHtml(title=sec["title"], file_name=f"chap_{i}.xhtml", lang="en")
        c.content = f"<h1>{sec['title']}</h1>" + "".join(f"<p>{p}</p>" for p in pages)
        book.add_item(c)
        items.append(c)
    book.toc = tuple(items)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", *items]
    epub.write_epub(out, book)


def write_xml(out: Path, sections, pages):
    body = "".join(f"<section title=\"{sec['title']}\"/>" for sec in sections)
    text = "".join(f"<p>{p}</p>" for p in pages)
    out.write_text(f"<doc>{body}<content>{text}</content></doc>", encoding="utf-8")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pdf", type=Path, nargs="?", default=Path("data/input/pipeline/BHT_CV32A65X_with_requirements_noannots.pdf"))
    ap.add_argument("--out", type=Path, default=Path("data/input/parity_gen"))
    args = ap.parse_args()

    pdf = args.pdf
    stem = pdf.stem
    target = args.out / stem
    target.mkdir(parents=True, exist_ok=True)

    sections = extract_sections(pdf)
    pages = extract_body(pdf)

    write_md(target / f"{stem}.md", sections, pages)
    write_html(target / f"{stem}.html", sections, pages)
    write_docx(target / f"{stem}.docx", sections, pages)
    write_rts(target / f"{stem}.rst", sections, pages)
    write_pptx(target / f"{stem}.pptx", sections, pages)
    write_xlsx(target / f"{stem}.xlsx", pages)
    write_epub(target / f"{stem}.epub", sections, pages)
    write_xml(target / f"{stem}.xml", sections, pages)

    print(f"Wrote parity fixtures to {target}")


if __name__ == "__main__":
    main()
