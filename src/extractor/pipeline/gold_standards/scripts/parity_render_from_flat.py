#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = [
#   "python-docx>=1.1.2",
#   "python-pptx>=0.6.23",
#   "openpyxl>=3.1.2",
#   "ebooklib>=0.18",
# ]
# ///
"""Render equivalent documents in multiple formats from a canonical flattened JSON.

Input: a flattened JSON (Stage 10 style) derived from the PDF. Output formats:
HTML, Markdown, RST, XML, DOCX, PPTX, XLSX, EPUB.
All outputs emit exactly one block per input block, preserving order.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import List, Dict, Any
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from ebooklib import epub


def load_blocks(flat_path: Path) -> List[Dict[str, Any]]:
    data = json.loads(flat_path.read_text())
    return data if isinstance(data, list) else []


def write_md(out: Path, blocks: List[Dict[str, Any]]):
    lines = []
    for blk in blocks:
        if blk.get("type") == "table" and blk.get("cells"):
            for row in blk["cells"]:
                vals = list(row.values()) if isinstance(row, dict) else row
                lines.append(" | ".join(str(v) for v in vals))
        else:
            lines.append(blk.get("text") or blk.get("content") or "")
    out.write_text("\n".join(lines), encoding="utf-8")


def write_html(out: Path, blocks: List[Dict[str, Any]]):
    parts = ["<html><body>"]
    for blk in blocks:
        if blk.get("type") == "table" and blk.get("cells"):
            parts.append("<table>")
            for row in blk["cells"]:
                vals = list(row.values()) if isinstance(row, dict) else row
                parts.append("<tr>" + "".join(f"<td>{v}</td>" for v in vals) + "</tr>")
            parts.append("</table>")
        else:
            parts.append(f"<p>{blk.get('text') or blk.get('content') or ''}</p>")
    parts.append("</body></html>")
    out.write_text("\n".join(parts), encoding="utf-8")


def write_rst(out: Path, blocks: List[Dict[str, Any]]):
    lines = []
    for blk in blocks:
        if blk.get("type") == "table" and blk.get("cells"):
            for row in blk["cells"]:
                vals = list(row.values()) if isinstance(row, dict) else row
                lines.append(" | ".join(str(v) for v in vals))
        else:
            lines.append(blk.get("text") or blk.get("content") or "")
    out.write_text("\n".join(lines), encoding="utf-8")


def write_xml(out: Path, blocks: List[Dict[str, Any]]):
    parts = ["<doc>"]
    for blk in blocks:
        if blk.get("type") == "table" and blk.get("cells"):
            parts.append("<table>")
            for row in blk["cells"]:
                vals = list(row.values()) if isinstance(row, dict) else row
                parts.append("<row>" + "".join(f"<col>{v}</col>" for v in vals) + "</row>")
            parts.append("</table>")
        else:
            parts.append(f"<p>{blk.get('text') or blk.get('content') or ''}</p>")
    parts.append("</doc>")
    out.write_text("\n".join(parts), encoding="utf-8")


def write_docx(out: Path, blocks: List[Dict[str, Any]]):
    doc = Document()
    for blk in blocks:
        if blk.get("type") == "table" and blk.get("cells"):
            rows = blk["cells"]
            first = rows[0]
            cols = len(first.values()) if isinstance(first, dict) else len(first)
            table = doc.add_table(rows=len(rows), cols=cols)
            for r_idx, row in enumerate(rows):
                vals = list(row.values()) if isinstance(row, dict) else row
                for c_idx, v in enumerate(vals):
                    table.cell(r_idx, c_idx).text = str(v)
        else:
            doc.add_paragraph(blk.get("text") or blk.get("content") or "")
    doc.save(out)


def write_pptx(out: Path, blocks: List[Dict[str, Any]]):
    prs = Presentation()
    bullet_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_layout)
    slide.shapes.title.text = "Content"
    tf = slide.shapes.placeholders[1].text_frame
    for blk in blocks:
        if blk.get("type") == "table" and blk.get("cells"):
            tf.add_paragraph().text = "[table]" + (blk.get("title") or "")
        else:
            tf.add_paragraph().text = blk.get("text") or blk.get("content") or ""
    prs.save(out)


def write_xlsx(out: Path, blocks: List[Dict[str, Any]]):
    wb = Workbook()
    ws = wb.active
    ws.title = "Content"
    row = 1
    for blk in blocks:
        if blk.get("type") == "table" and blk.get("cells"):
            for r in blk["cells"]:
                vals = list(r.values()) if isinstance(r, dict) else r
                for c_idx, v in enumerate(vals, start=1):
                    ws.cell(row=row, column=c_idx, value=v)
                row += 1
        else:
            ws.cell(row=row, column=1, value=blk.get("text") or blk.get("content") or "")
            row += 1
    wb.save(out)


def write_epub(out: Path, blocks: List[Dict[str, Any]]):
    book = epub.EpubBook()
    book.set_title(out.stem)
    book.set_language("en")
    items = []
    c = epub.EpubHtml(title="Content", file_name="chap_1.xhtml", lang="en")
    html = []
    for blk in blocks:
        if blk.get("type") == "table" and blk.get("cells"):
            rows = blk["cells"]
            html.append("<table>")
            for r in rows:
                vals = list(r.values()) if isinstance(r, dict) else r
                html.append("<tr>" + "".join(f"<td>{v}</td>" for v in vals) + "</tr>")
            html.append("</table>")
        else:
            html.append(f"<p>{blk.get('text') or blk.get('content') or ''}</p>")
    c.content = "".join(html)
    book.add_item(c)
    items.append(c)
    book.toc = tuple(items)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", *items]
    epub.write_epub(out, book)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--flat", type=Path, required=True, help="Flattened JSON from PDF")
    ap.add_argument("--out", type=Path, required=True, help="Output dir for rendered fixtures")
    args = ap.parse_args()

    blocks = load_blocks(args.flat)
    args.out.mkdir(parents=True, exist_ok=True)

    write_md(args.out / "reflow.md", blocks)
    write_html(args.out / "reflow.html", blocks)
    write_rst(args.out / "reflow.rst", blocks)
    write_xml(args.out / "reflow.xml", blocks)
    write_docx(args.out / "reflow.docx", blocks)
    write_pptx(args.out / "reflow.pptx", blocks)
    write_xlsx(args.out / "reflow.xlsx", blocks)
    write_epub(args.out / "reflow.epub", blocks)

    print(f"Rendered {len(blocks)} blocks to {args.out}")


if __name__ == "__main__":
    main()
