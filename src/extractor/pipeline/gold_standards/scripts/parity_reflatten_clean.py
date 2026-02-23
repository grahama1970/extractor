#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = []
# ///
"""Re-flatten clean artifacts (html/md) into a simple JSON block list.

This bypasses providers: it treats each paragraph/table row as one block in order.
Used only for deterministic parity against the canonical flattened PDF blocks.
"""
from __future__ import annotations
import argparse
import json
from pathlib import Path
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


def reflatten_html(path: Path):
    soup = BeautifulSoup(path.read_text(encoding="utf-8"), "html.parser")
    blocks = []
    for div in soup.find_all("div", class_="block"):
        if "table" in div.get("class", []):
            # one block per table, ignore rows for count parity
            blocks.append({"type": "table", "text": "table"})
        else:
            txt = div.get_text(strip=True)
            if txt:
                blocks.append({"type": "text", "text": txt})
    return blocks


def reflatten_md(path: Path):
    blocks = []
    for raw in path.read_text(encoding="utf-8").splitlines():
        line = raw.strip()
        if not line:
            continue
        if line.startswith("TABLE:"):
            blocks.append({"type": "table", "text": "table"})
        else:
            # drop the TEXT: prefix if present
            if line.startswith("TEXT:"):
                line = line[len("TEXT:") :].strip()
            blocks.append({"type": "text", "text": line})
    return blocks


def reflatten_rst(path: Path):
    blocks = []
    for raw in path.read_text(encoding="utf-8").splitlines():
        line = raw.strip()
        if not line:
            continue
        if line.startswith("TABLE:"):
            blocks.append({"type": "table", "text": "table"})
        else:
            if line.startswith("TEXT:"):
                line = line[len("TEXT:") :].strip()
            blocks.append({"type": "text", "text": line})
    return blocks


def reflatten_docx(path: Path):
    doc = Document(path)
    blocks = []
    for para in doc.paragraphs:
        txt = para.text.strip()
        if txt:
            blocks.append({"type": "text", "text": txt})
    for table in doc.tables:
        blocks.append({"type": "table", "text": "table"})
    return blocks


def reflatten_pptx(path: Path):
    prs = Presentation(path)
    blocks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "has_text_frame"):
                continue
            if shape.has_text_frame:
                txt = " ".join(p.text for p in shape.text_frame.paragraphs).strip()
                if txt:
                    blocks.append({"type": "text", "text": txt})
        # treat each table shape as one table block
        for shape in slide.shapes:
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                blocks.append({"type": "table", "text": "table"})
    return blocks


def reflatten_xlsx(path: Path):
    wb = load_workbook(path, read_only=True, data_only=True)
    blocks = []
    for ws in wb.worksheets:
        # heuristically treat each sheet as one table block if it has any non-empty cell
        has_data = any(cell.value not in (None, "") for row in ws.iter_rows() for cell in row)
        if has_data:
            blocks.append({"type": "table", "text": "table"})
    return blocks


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--html", type=Path, required=False)
    ap.add_argument("--md", type=Path, required=False)
    ap.add_argument("--rst", type=Path, required=False)
    ap.add_argument("--docx", type=Path, required=False)
    ap.add_argument("--pptx", type=Path, required=False)
    ap.add_argument("--xlsx", type=Path, required=False)
    ap.add_argument("--out", type=Path, required=True)
    args = ap.parse_args()

    blocks = []
    if args.html:
        blocks = reflatten_html(args.html)
    elif args.md:
        blocks = reflatten_md(args.md)
    elif args.rst:
        blocks = reflatten_rst(args.rst)
    elif args.docx:
        blocks = reflatten_docx(args.docx)
    elif args.pptx:
        blocks = reflatten_pptx(args.pptx)
    elif args.xlsx:
        blocks = reflatten_xlsx(args.xlsx)

    args.out.write_text(json.dumps(blocks, indent=2), encoding="utf-8")
    print(f"Reflattened {len(blocks)} blocks → {args.out}")


if __name__ == "__main__":
    main()
