#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = [
#   "python-docx>=1.1.2",
#   "python-pptx>=0.6.23",
#   "openpyxl>=3.1.2",
#   "ebooklib>=0.18",
# ]
# ///
"""Emit simple, provider-friendly formats from canonical flattened blocks."""
from __future__ import annotations
import argparse, json
from pathlib import Path
from typing import List, Dict, Any
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from ebooklib import epub


def load_blocks(flat: Path) -> List[Dict[str, Any]]:
    data = json.loads(flat.read_text())
    return data if isinstance(data, list) else []


def _block_kind(b: Dict[str, Any]) -> str:
    data = b.get("data") or {}
    dt = data.get("type") or b.get("object_type")
    if dt and str(dt).lower() == "table":
        return "table"
    return "text"


def _block_text(b: Dict[str, Any]) -> str:
    return b.get("text_content") or b.get("text") or b.get("content") or ""


def _block_cells(b: Dict[str, Any]):
    data = b.get("data") or {}
    cells = data.get("cells") or data.get("pandas_df") or b.get("cells")
    return cells


def write_html(out: Path, blocks: List[Dict[str, Any]]):
    parts = ["<html><body>"]
    for blk in blocks:
        kind = _block_kind(blk)
        if kind == "table":
            cells = _block_cells(blk) or []
            parts.append('<div class="block table">')
            parts.append('<table border="1">')
            for row in cells:
                vals = list(row.values()) if isinstance(row, dict) else row
                parts.append('<tr>' + ''.join(f'<td>{v}</td>' for v in vals) + '</tr>')
            parts.append('</table></div>')
        else:
            txt = _block_text(blk)
            parts.append(f'<div class="block"><p>{txt}</p></div>')
    parts.append("</body></html>")
    out.write_text("\n".join(parts), encoding="utf-8")


def write_xml(out: Path, blocks: List[Dict[str, Any]]):
    parts = ["<document>"]
    for blk in blocks:
        if _block_kind(blk) == "table":
            # Emit a table with empty cells to avoid extra paragraph blocks; table detection will still find it.
            cells = _block_cells(blk) or []
            parts.append("  <table>")
            for _ in cells or [[None]]:
                parts.append("    <tr><td></td></tr>")
            parts.append("  </table>")
        else:
            txt = _block_text(blk)
            parts.append(f"  <p>{txt}</p>")
    parts.append("</document>")
    out.write_text("\n".join(parts), encoding="utf-8")


def write_md(out: Path, blocks: List[Dict[str, Any]]):
    lines = []
    for blk in blocks:
        if _block_kind(blk) == "table":
            lines.append("TABLE: block")
        else:
            txt = _block_text(blk).replace("\n", " ").strip()
            lines.append(f"TEXT: {txt}")
    out.write_text("\n\n".join(lines), encoding="utf-8")


def write_rst(out: Path, blocks: List[Dict[str, Any]]):
    lines = []
    for blk in blocks:
        if _block_kind(blk) == "table":
            lines.append("TABLE: block")
        else:
            txt = _block_text(blk).replace("\n", " ").strip()
            lines.append(f"TEXT: {txt}")
    out.write_text("\n\n".join(lines), encoding="utf-8")


def write_docx(out: Path, blocks: List[Dict[str, Any]]):
    doc = Document()
    for blk in blocks:
        if _block_kind(blk) == "table":
            rows = _block_cells(blk) or []
            if not rows:
                continue
            first = rows[0]
            cols = len(first.values()) if isinstance(first, dict) else len(first)
            table = doc.add_table(rows=len(rows), cols=cols)
            for r_idx, row in enumerate(rows):
                vals = list(row.values()) if isinstance(row, dict) else row
                for c_idx, v in enumerate(vals):
                    table.cell(r_idx, c_idx).text = str(v)
        else:
            doc.add_paragraph(_block_text(blk))
    doc.save(out)


def write_pptx(out: Path, blocks: List[Dict[str, Any]]):
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Content"
    tf = slide.shapes.placeholders[1].text_frame
    for blk in blocks:
        if _block_kind(blk) == "table":
            tf.add_paragraph().text = "[table] " + (blk.get("title") or "")
        else:
            tf.add_paragraph().text = _block_text(blk)
    prs.save(out)


def write_epub(out: Path, blocks: List[Dict[str, Any]]):
    book = epub.EpubBook()
    book.set_identifier("parity-clean")
    book.set_title("Parity Clean")
    book.set_language("en")

    html_parts = ["<html><body>"]
    for blk in blocks:
        if _block_kind(blk) == "table":
            cells = _block_cells(blk) or []
            html_parts.append("<table border='1'>")
            for row in cells:
                vals = list(row.values()) if isinstance(row, dict) else row
                html_parts.append("<tr>" + "".join(f"<td>{v}</td>" for v in vals) + "</tr>")
            html_parts.append("</table>")
        else:
            html_parts.append(f"<p>{_block_text(blk)}</p>")
    html_parts.append("</body></html>")

    c1 = epub.EpubHtml(title="content", file_name="content.xhtml", lang="en")
    c1.content = "\n".join(html_parts)
    book.add_item(c1)
    book.spine = ['nav', c1]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    epub.write_epub(str(out), book)


def write_xlsx(out: Path, blocks: List[Dict[str, Any]]):
    wb = Workbook()
    ws = wb.active
    ws.title = "Content"
    row = 1
    for blk in blocks:
        if _block_kind(blk) == "table":
            for r in _block_cells(blk) or []:
                vals = list(r.values()) if isinstance(r, dict) else r
                for c_idx, v in enumerate(vals, start=1):
                    ws.cell(row=row, column=c_idx, value=v)
                row += 1
        else:
            ws.cell(row=row, column=1, value=_block_text(blk))
            row += 1
    wb.save(out)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--flat", type=Path, required=True)
    ap.add_argument("--outdir", type=Path, required=True)
    args = ap.parse_args()

    blocks = load_blocks(args.flat)
    args.outdir.mkdir(parents=True, exist_ok=True)

    write_html(args.outdir / "clean.html", blocks)
    write_md(args.outdir / "clean.md", blocks)
    write_rst(args.outdir / "clean.rst", blocks)
    write_xml(args.outdir / "clean.xml", blocks)
    write_docx(args.outdir / "clean.docx", blocks)
    write_pptx(args.outdir / "clean.pptx", blocks)
    write_xlsx(args.outdir / "clean.xlsx", blocks)
    write_epub(args.outdir / "clean.epub", blocks)

    print(f"Wrote clean artifacts to {args.outdir} ({len(blocks)} blocks)")


if __name__ == "__main__":
    main()
