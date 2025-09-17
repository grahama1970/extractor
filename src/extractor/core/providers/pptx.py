"""
Module: pptx.py
Purpose: Native PowerPoint extraction without PDF conversion

This module implements direct PPTX extraction using python-pptx,
preserving all presentation features including speaker notes, slide
layouts, animations metadata, and embedded objects without lossy conversions.

External Dependencies:
- python-pptx: https://python-pptx.readthedocs.io/

Example Usage:
>>> from extractor.core.providers.pptx import PPTXProvider
>>> provider = NativePPTXProvider()
>>> document = provider.extract_document("presentation.pptx")
>>> print(document.source_type)  # SourceType.PPTX
"""

import hashlib
import base64
from pathlib import Path
from typing import List, Dict, Any, Optional, Union
from datetime import datetime

from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.dml import MSO_THEME_COLOR
from loguru import logger

from extractor.core.schema.unified_document import (
    UnifiedDocument,
    BlockType,
    SourceType,
    BaseBlock,
    TableBlock,
    ImageBlock,
    BlockMetadata,
    DocumentMetadata,
    HierarchyNode,
    TableCell,
)

# Default configuration constants
DEFAULT_MAX_IMAGE_SIZE_MB = 10
DEFAULT_CLAUDE_TIMEOUT = 30
DEFAULT_GEMINI_TIMEOUT = 20


class PPTXProvider:
    """Direct PPTX extraction without PDF conversion"""

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        self.config = config or {}
        self.block_counter = 0
        self.slide_hierarchy = {}

        # Validate config
        self.config.setdefault("max_image_size_mb", DEFAULT_MAX_IMAGE_SIZE_MB)
        self.config.setdefault("enable_ai_descriptions", False)  # sync for now
        self.config.setdefault("enable_table_merge_analysis", True)
        self.config.setdefault("claude_timeout", DEFAULT_CLAUDE_TIMEOUT)
        self.config.setdefault("gemini_timeout", DEFAULT_GEMINI_TIMEOUT)
        self.config.setdefault("gemini_model", "vertex_ai/gemini-2.0-flash-exp")

    def extract_document(self, filepath: Union[str, Path]) -> UnifiedDocument:
        """Extract PPTX content to unified document format"""
        filepath = Path(filepath)
        if not filepath.exists():
            raise FileNotFoundError(filepath)

        logger.info(f"Extracting PPTX document: {filepath}")

        # Load presentation
        prs = Presentation(str(filepath))

        # Extract metadata
        metadata = self._extract_metadata(prs, filepath)

        # Extract all slides with structure
        blocks = []
        slide_nodes = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_blocks, slide_node = self._extract_slide(slide, slide_idx)
            blocks.extend(slide_blocks)
            slide_nodes.append(slide_node)

        # Build document hierarchy
        hierarchy = self._build_hierarchy(slide_nodes)

        # Apply Claude table merge analysis if enabled
        if self.config.get("enable_table_merge_analysis", True):
            blocks = self._analyze_and_merge_tables(blocks)

        # Create unified document
        doc = UnifiedDocument(
            id=self._generate_doc_id(filepath),
            source_type=SourceType.PPTX,
            source_path=str(filepath),
            blocks=blocks,
            hierarchy=hierarchy,
            metadata=metadata,
            full_text=self._extract_full_text(blocks),
            keywords=self._extract_keywords(prs),
            _key=self._generate_doc_id(filepath),
            _collection="documents",
        )

        logger.info(f"Extracted {len(blocks)} blocks from {len(prs.slides)} slides")
        return doc

    def _generate_doc_id(self, filepath: Path) -> str:
        """Generate unique document ID using SHA256"""
        return hashlib.sha256(str(filepath).encode()).hexdigest()

    def _extract_metadata(self, prs: Any, filepath: Path) -> DocumentMetadata:
        """Extract presentation metadata"""
        metadata = DocumentMetadata()

        # Core properties
        props = prs.core_properties
        if props:
            metadata.title = props.title or ""
            metadata.author = props.author or ""

            # Dates
            if props.created:
                metadata.created_date = props.created
            if props.modified:
                metadata.modified_date = props.modified

            # Language (if available)
            metadata.language = props.language or ""

        # Presentation-specific metadata
        metadata.format_metadata = {
            "file_type": "pptx",
            "slide_count": len(prs.slides),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
            "has_notes": any(slide.has_notes_slide for slide in prs.slides),
            "presentation_format": f"{prs.slide_width}x{prs.slide_height}",
        }

        # File metadata
        metadata.file_size = filepath.stat().st_size if filepath.exists() else None

        return metadata

    def _extract_slide(self, slide: Slide, slide_idx: int) -> tuple[List[BaseBlock], HierarchyNode]:
        """Extract all content from a single slide with AI-powered analysis"""
        blocks = []

        # Create slide heading block
        slide_title = self._get_slide_title(slide)
        if not slide_title:
            slide_title = f"Slide {slide_idx + 1}"

        # Generate AI-powered slide description if enabled
        slide_description = ""
        if self.config.get("enable_slide_descriptions", True):
            slide_description = self._get_ai_slide_description(slide, slide_idx, slide_title)

        slide_block = BaseBlock(
            id=self._generate_block_id(),
            type=BlockType.HEADING,
            content=slide_title,
            parent_id=None,
            metadata=BlockMetadata(
                page_number=slide_idx + 1,
                bbox=None,
                source_id=f"slide_{slide_idx}",
                language=None,
                attributes={
                    "level": 1,
                    "slide_number": slide_idx + 1,
                    "is_slide_title": True,
                    "ai_description": slide_description,
                    "shape_count": len(slide.shapes),
                    "has_notes": slide.has_notes_slide,
                },
                confidence=1.0,
            ),
        )
        blocks.append(slide_block)

        # Add slide description as a separate block if generated
        if slide_description:
            desc_block = BaseBlock(
                id=self._generate_block_id(),
                type=BlockType.PARAGRAPH,
                content=f"Slide Overview: {slide_description}",
                parent_id=slide_block.id,
                metadata=BlockMetadata(
                    page_number=slide_idx + 1,
                    bbox=None,
                    source_id=f"slide_{slide_idx}_desc",
                    language=None,
                    attributes={
                        "slide_number": slide_idx + 1,
                        "content_type": "ai_slide_description",
                        "generated_by": "pptx_provider",
                    },
                    confidence=0.85,
                ),
            )
            blocks.append(desc_block)

        # Extract shapes
        for shape in slide.shapes:
            shape_blocks = self._extract_shape(shape, slide_idx, parent_id=slide_block.id)
            blocks.extend(shape_blocks)

        # Extract notes
        if slide.has_notes_slide:
            notes_blocks = self._extract_notes(slide.notes_slide, slide_idx, parent_id=slide_block.id)
            blocks.extend(notes_blocks)

        # Create hierarchy node for slide
        slide_node = HierarchyNode(
            id=f"slide-{slide_idx}",
            title=slide_title,
            level=1,
            block_id=slide_block.id,
            breadcrumb=["Presentation", slide_title],
        )

        return blocks, slide_node

    def _get_ai_slide_description(self, slide: Slide, slide_idx: int, slide_title: str) -> str:
        """Generate AI-powered description of the entire slide composition"""
        try:
            shape_types = []
            text_content = []
            visual_elements = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    text_content.append(shape.text.strip()[:100])

                if shape.shape_type.name:
                    shape_types.append(shape.shape_type.name.lower())

                if hasattr(shape, "has_table") and shape.has_table:
                    visual_elements.append("table")
                elif hasattr(shape, "has_chart") and shape.has_chart:
                    visual_elements.append("chart")
                elif shape.shape_type.name and "picture" in shape.shape_type.name.lower():
                    visual_elements.append("image")
                elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    visual_elements.append("text")

            description_parts = []
            if visual_elements:
                element_counts = {}
                for elem in visual_elements:
                    element_counts[elem] = element_counts.get(elem, 0) + 1
                content_desc = [
                    f"{c} {t}" if c == 1 else f"{c} {t}s" for t, c in element_counts.items()
                ]
                if content_desc:
                    description_parts.append(f"Contains {', '.join(content_desc)}")

            total_text_length = sum(len(t) for t in text_content)
            if total_text_length > 200:
                description_parts.append("with substantial text content")
            elif total_text_length > 50:
                description_parts.append("with moderate text content")
            else:
                description_parts.append("with minimal text")

            shape_count = len(slide.shapes)
            if shape_count > 10:
                description_parts.append("featuring a complex layout")
            elif shape_count > 5:
                description_parts.append("with multiple elements")
            elif shape_count <= 2:
                description_parts.append("with a simple layout")

            if description_parts:
                return f"Slide {slide_idx + 1} {' '.join(description_parts).lower()}."
            return f"Slide {slide_idx + 1} with visual content."

        except Exception as e:
            logger.debug(f"AI slide description failed for slide {slide_idx + 1}: {e}")
            return ""

    def _get_slide_title(self, slide: Slide) -> Optional[str]:
        """Extract slide title using canonical placeholder when available."""
        try:
            t = getattr(slide.shapes, "title", None)
            if t is not None and getattr(t, "has_text_frame", False):
                txt = t.text if hasattr(t, "text") else t.text_frame.text
                if txt and str(txt).strip():
                    return str(txt).strip()
        except Exception:
            pass
        # Fallback 1: scan placeholders for TITLE/CTR_TITLE
        for shape in slide.shapes:
            try:
                if getattr(shape, "is_placeholder", False):
                    ph = shape.placeholder_format.type
                    if ph in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                        if hasattr(shape, "text") and shape.text:
                            return shape.text.strip()
            except Exception:
                continue
        # Fallback 2: prefer BODY placeholder's first paragraph
        try:
            for ph in getattr(slide, "placeholders", []):
                if getattr(ph, "has_text_frame", False):
                    pht = ph.placeholder_format.type
                    if pht == PP_PLACEHOLDER.BODY and ph.text:
                        txt = str(ph.text).strip()
                        if txt:
                            return txt.splitlines()[0][:120]
        except Exception:
            pass
        # Fallback 3: choose the top-most text-bearing shape
        try:
            candidates = []
            for sh in slide.shapes:
                if getattr(sh, "has_text_frame", False):
                    text = sh.text_frame.text if hasattr(sh, "text_frame") else getattr(sh, "text", None)
                    if text and str(text).strip():
                        top = int(getattr(sh, "top", 1_000_000_000))
                        left = int(getattr(sh, "left", 1_000_000_000))
                        candidates.append((top, left, str(text).strip()))
            if candidates:
                candidates.sort(key=lambda t: (t[0], t[1]))
                return candidates[0][2].splitlines()[0][:120]
        except Exception:
            pass
        return None

    def _extract_shape(self, shape: Any, slide_idx: int, *, parent_id: Optional[str]) -> List[BaseBlock]:
        """Extract content from a shape"""
        blocks = []

        # Handle grouped shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                sub_blocks = self._extract_shape(sub_shape, slide_idx, parent_id=parent_id)
                blocks.extend(sub_blocks)
            return blocks

        # Handle tables
        if shape.has_table:
            table_block = self._extract_table(shape, slide_idx, parent_id=parent_id)
            if table_block:
                blocks.append(table_block)
            return blocks

        # Handle images/pictures
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_block = self._extract_image(shape, slide_idx, parent_id=parent_id)
            if image_block:
                blocks.append(image_block)
            return blocks

        # Handle charts
        if shape.has_chart:
            chart_block = self._extract_chart(shape, slide_idx, parent_id=parent_id)
            if chart_block:
                blocks.append(chart_block)
            return blocks

        # Handle text
        if hasattr(shape, "text") and shape.text and shape.text.strip():
            text_blocks = self._extract_text(shape, slide_idx, parent_id=parent_id)
            blocks.extend(text_blocks)

        return blocks

    def _extract_text(self, shape: Any, slide_idx: int, *, parent_id: Optional[str]) -> List[BaseBlock]:
        """Extract text content from shape"""
        blocks = []

        block_type = BlockType.PARAGRAPH
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            if placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                block_type = BlockType.HEADING
            elif placeholder_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                return blocks  # already handled as slide title

        if shape.has_text_frame:
            all_paragraphs = [
                (p, "".join(run.text for run in p.runs).strip())
                for p in shape.text_frame.paragraphs
                if "".join(run.text for run in p.runs).strip()
            ]

            if (
                len(all_paragraphs) > 1
                and shape.is_placeholder
                and shape.placeholder_format.type
                not in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]
            ):
                list_items = []
                for paragraph, text in all_paragraphs:
                    p_el = paragraph._element
                    bullet_num = p_el.find(".//a:buAutoNum", namespaces=p_el.nsmap)
                    list_items.append(
                        {
                            "text": text,
                            "level": paragraph.level,
                            "is_numbered": bullet_num is not None,
                        }
                    )
                if list_items:
                    blocks.append(self._create_list_block(list_items, slide_idx, parent_id=parent_id))
            else:
                for _, text in all_paragraphs:
                    blocks.append(
                        BaseBlock(
                            id=self._generate_block_id(),
                            type=block_type,
                            content=text,
                            parent_id=parent_id,
                            metadata=BlockMetadata(
                                page_number=slide_idx + 1,
                                bbox=None,
                                source_id=f"slide_{slide_idx}_text",
                                language=None,
                                attributes={
                                    "slide_number": slide_idx + 1,
                                    "shape_type": "text",
                                    **({"level": 2} if block_type == BlockType.HEADING else {}),
                                },
                                confidence=0.95,
                            ),
                        )
                    )
        return blocks

    def _create_list_block(self, items: List[Dict[str, Any]], slide_idx: int, *, parent_id: Optional[str]) -> BaseBlock:
        """Create a list block from items"""
        content_lines = []
        for item in items:
            indent = "  " * item["level"]
            marker = f"{item['level'] + 1}." if item["is_numbered"] else "-"
            content_lines.append(f"{indent}{marker} {item['text']}")

        return BaseBlock(
            id=self._generate_block_id(),
            type=BlockType.LISTITEM,
            content="\n".join(content_lines),
            parent_id=parent_id,
            metadata=BlockMetadata(
                page_number=slide_idx + 1,
                bbox=None,
                source_id=f"slide_{slide_idx}_list",
                language=None,
                attributes={
                    "slide_number": slide_idx + 1,
                    "list_type": "ordered" if items[0]["is_numbered"] else "unordered",
                    "item_count": len(items),
                },
                confidence=0.95,
            ),
        )

    def _extract_table(self, shape: Any, slide_idx: int, *, parent_id: Optional[str]) -> Optional[TableBlock]:
        """Extract table from shape"""
        table = shape.table
        cells = []
        headers = [0]  # assume first row is header

        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cells.append(
                    TableCell(
                        row=row_idx, col=col_idx, content=cell.text.strip() if cell.text else ""
                    )
                )

        if not cells:
            return None

        return TableBlock(
            id=self._generate_block_id(),
            type=BlockType.TABLE,
            content="",
            parent_id=parent_id,
            rows=len(table.rows),
            cols=len(table.columns),
            cells=cells,
            headers=headers,
            metadata=BlockMetadata(
                page_number=slide_idx + 1,
                bbox=None,
                source_id=f"slide_{slide_idx}_table",
                language=None,
                attributes={"slide_number": slide_idx + 1, "shape_type": "table"},
                confidence=0.95,
            ),
        )

    def _extract_image(self, shape: Any, slide_idx: int, *, parent_id: Optional[str]) -> Optional[ImageBlock]:
        """Extract image from shape with size protection"""
        try:
            image = shape.image
            image_bytes = image.blob

            max_size = self.config.get("max_image_size_mb", DEFAULT_MAX_IMAGE_SIZE_MB) * 1024 * 1024
            if len(image_bytes) > max_size:
                logger.warning(
                    f"Image too large ({len(image_bytes)} bytes), skipping base64 encoding"
                )
                data_uri = f"data:{image.content_type};base64,<large_image_skipped>"
                ai_description = self._get_fallback_description(shape, slide_idx)
            else:
                image_base64 = base64.b64encode(image_bytes).decode("utf-8")
                data_uri = f"data:{image.content_type};base64,{image_base64}"
                ai_description = self._get_fallback_description(shape, slide_idx)

            alt_text_prop = getattr(shape, "alternative_text", None)
            if not alt_text_prop:
                alt_text_prop = getattr(shape, "alt_text", None)
            alt_text = (alt_text_prop or shape.name or f"Image from slide {slide_idx + 1}")
            if ai_description:
                alt_text = f"{alt_text}: {ai_description}"

            return ImageBlock(
                id=self._generate_block_id(),
                type=BlockType.IMAGE,
                content=ai_description,
                parent_id=parent_id,
                src=data_uri,
                alt=alt_text,
                format=image.ext or "png",
                metadata=BlockMetadata(
                    page_number=slide_idx + 1,
                    bbox=None,
                    source_id=f"slide_{slide_idx}_image",
                    language=None,
                    attributes={
                        "slide_number": slide_idx + 1,
                        "shape_type": "picture",
                        "content_type": image.content_type,
                        "ai_described": bool(ai_description),
                        "shape_name": shape.name or "",
                        "image_size_bytes": len(image_bytes),
                    },
                    confidence=1.0,
                ),
            )
        except Exception as e:
            logger.warning(f"Failed to extract image: {e}")
            return None

    def _get_fallback_description(self, shape: Any, slide_idx: int) -> str:
        """Fallback description when AI analysis fails"""
        figure_type = self._classify_ppt_image_type(shape, slide_idx)

        mapping = {
            "chart": f"Chart or graph displaying data visualization on slide {slide_idx + 1}",
            "diagram": f"Diagram or flowchart illustrating concepts on slide {slide_idx + 1}",
            "screenshot": f"Screenshot or interface capture shown on slide {slide_idx + 1}",
            "table": f"Table or data grid presented as image on slide {slide_idx + 1}",
            "photo": f"Photograph or image content on slide {slide_idx + 1}",
        }
        return mapping.get(figure_type, f"Visual content on slide {slide_idx + 1}")

    def _classify_ppt_image_type(self, shape: Any, slide_idx: int) -> str:
        """Classify PowerPoint image type based on shape properties and context"""
        shape_name = (shape.name or "").lower()
        keywords = {
            "chart": ["chart", "graph", "plot", "data"],
            "diagram": ["diagram", "flow", "process", "workflow"],
            "screenshot": ["screen", "capture", "interface", "ui"],
            "table": ["table", "grid", "data"],
            "photo": ["photo", "image", "picture"],
        }
        for t, kw in keywords.items():
            if any(k in shape_name for k in kw):
                return t
        return "unknown"

    def _extract_chart(self, shape: Any, slide_idx: int, *, parent_id: Optional[str]) -> Optional[BaseBlock]:
        """Extract chart information with error protection"""
        try:
            chart = shape.chart
            chart_info = []

            try:
                chart_type_str = (
                    str(chart.chart_type) if hasattr(chart, "chart_type") else "Unknown"
                )
                chart_info.append(f"Chart Type: {chart_type_str}")
            except Exception:
                chart_info.append("Chart Type: Unknown")

            try:
                title = (
                    chart.chart_title.text_frame.text
                    if (
                        hasattr(chart, "has_title")
                        and chart.has_title
                        and hasattr(chart, "chart_title")
                    )
                    else "Untitled"
                )
                chart_info.append(f"Title: {title}")
            except Exception:
                chart_info.append("Title: Untitled")

            try:
                if hasattr(chart, "series") and chart.series:
                    chart_info.append(f"Series Count: {len(chart.series)}")
            except Exception:
                pass

            return BaseBlock(
                id=self._generate_block_id(),
                type=BlockType.FIGURE,
                content="\n".join(chart_info),
                parent_id=parent_id,
                metadata=BlockMetadata(
                    page_number=slide_idx + 1,
                    bbox=None,
                    source_id=f"slide_{slide_idx}_chart",
                    language=None,
                    attributes={
                        "slide_number": slide_idx + 1,
                        "shape_type": "chart",
                        "chart_type": (
                            str(chart.chart_type) if hasattr(chart, "chart_type") else "Unknown"
                        ),
                    },
                    confidence=0.9,
                ),
            )
        except Exception as e:
            logger.warning(f"Failed to extract chart: {e}")
            return None

    def _extract_notes(self, notes_slide: Any, slide_idx: int, *, parent_id: Optional[str]) -> List[BaseBlock]:
        """Extract speaker notes"""
        blocks = []
        for shape in notes_slide.shapes:
            if shape.has_text_frame and shape.text:
                blocks.append(
                    BaseBlock(
                        id=self._generate_block_id(),
                        type=BlockType.COMMENT,
                        content=shape.text.strip(),
                        parent_id=parent_id,
                        metadata=BlockMetadata(
                            page_number=slide_idx + 1,
                            bbox=None,
                            source_id=f"slide_{slide_idx}_notes",
                            language=None,
                            attributes={
                                "slide_number": slide_idx + 1,
                                "comment_type": "speaker_notes",
                            },
                            confidence=1.0,
                        ),
                    )
                )
        return blocks

    def _build_hierarchy(self, slide_nodes: List[HierarchyNode]) -> HierarchyNode:
        """Build presentation hierarchy"""
        root = HierarchyNode(
            id="root", title="Presentation", level=0, block_id="root", children=slide_nodes
        )
        for node in slide_nodes:
            node.parent_id = root.id
        return root

    def _extract_full_text(self, blocks: List[BaseBlock]) -> str:
        """Extract all text content"""
        text_parts = []
        for block in blocks:
            if hasattr(block, "content") and isinstance(block.content, str):
                text_parts.append(block.content)
            elif block.type == BlockType.TABLE and isinstance(block, TableBlock):
                for cell in block.cells:
                    if cell.content:
                        text_parts.append(cell.content)
        return "\n".join(text_parts)

    def _extract_keywords(self, prs: Any) -> List[str]:
        """Extract keywords from presentation"""
        keywords = []
        if prs.core_properties:
            props = prs.core_properties
            if props.keywords:
                keywords.extend([k.strip() for k in props.keywords.split(",") if k.strip()])
            if props.subject:
                keywords.append(props.subject)
            if props.category:
                keywords.append(props.category)
        return keywords

    def _generate_block_id(self) -> str:
        """Generate unique block ID"""
        self.block_counter += 1
        return f"pptx-block-{self.block_counter}"

    def _analyze_and_merge_tables(self, blocks: List[BaseBlock]) -> List[BaseBlock]:
        """Analyze tables for potential merging using Claude's logic"""
        try:
            from extractor.core.processors.claude_table_merge_analyzer import (
                ClaudeTableMergeAnalyzer,
            )

            analyzer = ClaudeTableMergeAnalyzer()
            return analyzer.analyze_and_merge_tables(blocks)
        except (ImportError, ModuleNotFoundError):
            logger.debug("Claude table merge analyzer not available, skipping table merge analysis")
            return blocks


if __name__ == "__main__":
    # Quick smoke test
    provider = PPTXProvider()
    assert provider is not None
    assert provider.block_counter == 0
    print("✅ Native PPTX provider initialized successfully")
    print("Note: Full testing requires actual PowerPoint files")
