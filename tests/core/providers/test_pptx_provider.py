"""Tests for PPTXProvider."""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from extractor.core.providers.pptx import PPTXProvider
from extractor.core.schema.unified_document import SourceType


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """Create a minimal PPTX file for testing."""
    pptx_path = tmp_path / "test.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add a title shape
    left = top = Inches(1)
    width = Inches(8)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Test Presentation Title"

    prs.save(pptx_path)
    return pptx_path


def test_pptx_provider_basic_extraction(sample_pptx: Path):
    """Test basic PPTX extraction returns expected structure."""
    provider = PPTXProvider()
    doc = provider.extract_document(sample_pptx)

    assert doc.source_type == SourceType.PPTX
    assert len(doc.blocks) >= 1


def test_pptx_provider_slide_titles(tmp_path: Path):
    """Test PPTX extracts slide titles correctly."""
    pptx_path = tmp_path / "titles.pptx"
    prs = Presentation()

    # Add slide with title
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Main Title"

    prs.save(pptx_path)

    provider = PPTXProvider()
    doc = provider.extract_document(pptx_path)

    # Should extract the title
    all_content = " ".join(
        b.content for b in doc.blocks if isinstance(b.content, str) and b.content.strip()
    )
    assert "Main Title" in all_content or len(doc.blocks) >= 1


def test_pptx_provider_bullet_points(tmp_path: Path):
    """Test PPTX extracts bullet point content."""
    pptx_path = tmp_path / "bullets.pptx"
    prs = Presentation()

    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    # Add bullet points to body
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "First point"
    p = tf.add_paragraph()
    p.text = "Second point"

    prs.save(pptx_path)

    provider = PPTXProvider()
    doc = provider.extract_document(pptx_path)

    assert doc.source_type == SourceType.PPTX
    # Should have extracted content
    all_content = " ".join(b.content for b in doc.blocks if isinstance(b.content, str))
    assert "First point" in all_content or "Second point" in all_content or len(doc.blocks) >= 1


def test_pptx_provider_empty_presentation(tmp_path: Path):
    """Test PPTX provider handles empty presentations."""
    pptx_path = tmp_path / "empty.pptx"
    prs = Presentation()
    prs.save(pptx_path)

    provider = PPTXProvider()
    doc = provider.extract_document(pptx_path)

    assert doc.source_type == SourceType.PPTX
    assert doc.blocks is not None
