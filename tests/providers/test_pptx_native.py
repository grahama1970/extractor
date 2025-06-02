"""
Test native PowerPoint extraction without PDF conversion
"""

import pytest
import tempfile
import time
from pathlib import Path
from datetime import datetime

from marker.core.providers.pptx_native import NativePPTXProvider
from marker.core.schema.unified_document import (
    UnifiedDocument, BlockType, SourceType, BaseBlock, TableBlock,
    ImageBlock
)


class TestNativePPTXExtractor:
    """Test native PPTX extraction functionality"""
    
    def test_simple_presentation(self):
        """Test extraction of simple PPTX with text and slides"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            pytest.skip("python-pptx not available")
        
        # Create test presentation
        prs = Presentation()
        
        # Slide 1: Title slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide1.shapes.title
        subtitle = slide1.placeholders[1]
        title.text = "Test Presentation"
        subtitle.text = "Automated Testing"
        
        # Add metadata
        prs.core_properties.title = "Test Presentation" 
        prs.core_properties.author = "Test Author"
        prs.core_properties.subject = "Testing"
        prs.core_properties.keywords = "test, pptx, extraction"
        
        # Slide 2: Content slide
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        title2 = slide2.shapes.title
        title2.text = "Content Slide"
        
        # Add bullet points
        content = slide2.placeholders[1]
        tf = content.text_frame
        tf.text = "First bullet point"
        
        p = tf.add_paragraph()
        p.text = "Second bullet point"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = "Sub-bullet"
        p.level = 1
        
        # Save to temp file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            temp_path = Path(f.name)
        
        try:
            # Extract
            start_time = time.time()
            provider = NativePPTXProvider()
            extracted_doc = provider.extract_document(temp_path)
            duration = time.time() - start_time
            
            # Verify extraction
            assert extracted_doc.source_type == SourceType.PPTX
            assert extracted_doc.metadata.title == "Test Presentation"
            assert extracted_doc.metadata.author == "Test Author"
            
            # Check metadata
            assert extracted_doc.metadata.format_metadata['file_type'] == 'pptx'
            assert extracted_doc.metadata.format_metadata['slide_count'] == 2
            
            # Check keywords
            assert len(extracted_doc.keywords) >= 2
            
            # Check blocks - should have slide headings
            heading_blocks = [b for b in extracted_doc.blocks if b.type == BlockType.HEADING]
            assert len(heading_blocks) >= 3  # 2 slide titles + content slide title
            
            # Check hierarchy
            assert extracted_doc.hierarchy is not None
            assert extracted_doc.hierarchy.title == "Presentation"
            assert len(extracted_doc.hierarchy.children) == 2  # 2 slides
            
            # Check bullet list extraction
            list_blocks = [b for b in extracted_doc.blocks if b.type == BlockType.LISTITEM]
            assert len(list_blocks) >= 1
            
            # Verify list content
            list_content = list_blocks[0].content
            assert "First bullet point" in list_content
            assert "Second bullet point" in list_content
            assert "Sub-bullet" in list_content
            
            # Duration check - should be fast
            assert duration < 1.0
            
        finally:
            temp_path.unlink()
    
    def test_presentation_with_tables(self):
        """Test extraction of PPTX with tables"""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            pytest.skip("python-pptx not available")
        
        prs = Presentation()
        
        # Add slide with table
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        
        # Add table
        rows, cols = 3, 3
        left = Inches(1)
        top = Inches(2)
        width = Inches(6)
        height = Inches(2)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Header row
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        table.cell(0, 2).text = "Status"
        
        # Data rows
        table.cell(1, 0).text = "Test Item 1"
        table.cell(1, 1).text = "100"
        table.cell(1, 2).text = "Active"
        
        table.cell(2, 0).text = "Test Item 2"
        table.cell(2, 1).text = "200"
        table.cell(2, 2).text = "Pending"
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            temp_path = Path(f.name)
        
        try:
            provider = NativePPTXProvider()
            extracted_doc = provider.extract_document(temp_path)
            
            # Check tables
            table_blocks = [b for b in extracted_doc.blocks if b.type == BlockType.TABLE]
            assert len(table_blocks) >= 1
            
            table = table_blocks[0]
            assert table.rows == 3
            assert table.cols == 3
            
            # Check table content
            table_texts = [cell.content for cell in table.cells]
            assert "Test Item 1" in table_texts
            assert "100" in table_texts
            assert "Active" in table_texts
            
        finally:
            temp_path.unlink()
    
    def test_speaker_notes(self):
        """Test extraction of speaker notes"""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not available")
        
        prs = Presentation()
        
        # Add slide with notes
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Slide with Notes"
        
        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "These are important speaker notes for the presentation."
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            temp_path = Path(f.name)
        
        try:
            provider = NativePPTXProvider()
            extracted_doc = provider.extract_document(temp_path)
            
            # Check for speaker notes
            comment_blocks = [b for b in extracted_doc.blocks if b.type == BlockType.COMMENT]
            assert len(comment_blocks) >= 1
            
            # Verify notes content
            notes_found = False
            for block in comment_blocks:
                if block.metadata.attributes.get('comment_type') == 'speaker_notes':
                    assert "important speaker notes" in block.content
                    notes_found = True
                    break
            
            assert notes_found, "Speaker notes not found"
            
            # Check metadata
            assert extracted_doc.metadata.format_metadata['has_notes'] == True
            
        finally:
            temp_path.unlink()
    
    def test_images_extraction(self):
        """Test extraction of images from slides"""
        try:
            from pptx import Presentation
            from pptx.util import Inches
            from PIL import Image
            import io
        except ImportError:
            pytest.skip("python-pptx or PIL not available")
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Create a simple image
        img = Image.new('RGB', (100, 100), color='red')
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        
        # Add image to slide
        slide.shapes.add_picture(img_bytes, Inches(1), Inches(1), width=Inches(2))
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            temp_path = Path(f.name)
        
        try:
            provider = NativePPTXProvider()
            extracted_doc = provider.extract_document(temp_path)
            
            # Check for images
            image_blocks = [b for b in extracted_doc.blocks if b.type == BlockType.IMAGE]
            assert len(image_blocks) >= 1
            
            # Verify image is base64 encoded
            image = image_blocks[0]
            assert image.src.startswith("data:image/")
            assert "base64," in image.src
            
        finally:
            temp_path.unlink()
    
    def test_no_pdf_conversion(self):
        """HONEYPOT: Ensure no PDF conversion happens"""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not available")
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Direct Extraction Test"
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            temp_path = Path(f.name)
        
        try:
            provider = NativePPTXProvider()
            
            # Extraction should be very fast (no conversion overhead)
            start_time = time.time()
            extracted_doc = provider.extract_document(temp_path)
            duration = time.time() - start_time
            
            # Should be much faster than PDF conversion
            assert duration < 0.5
            
            # Check that we have PPTX-specific metadata
            assert extracted_doc.source_type == SourceType.PPTX
            assert 'slide_count' in extracted_doc.metadata.format_metadata
            assert 'slide_width' in extracted_doc.metadata.format_metadata
            
            # No PDF artifacts
            full_text = extracted_doc.full_text.lower() if extracted_doc.full_text else ""
            assert 'pdf' not in str(extracted_doc.metadata).lower()
            
        finally:
            temp_path.unlink()
    
    def test_complex_shapes(self):
        """Test extraction of grouped shapes and charts"""
        try:
            from pptx import Presentation
            from pptx.util import Inches
            from pptx.chart.data import ChartData
            from pptx.enum.chart import XL_CHART_TYPE
        except ImportError:
            pytest.skip("python-pptx not available")
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Add a chart
        chart_data = ChartData()
        chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
        chart_data.add_series('Sales', (30, 25, 40, 35))
        
        x, y, cx, cy = Inches(1), Inches(1), Inches(6), Inches(4)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        chart.has_title = True
        chart.chart_title.text_frame.text = "Quarterly Sales"
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            temp_path = Path(f.name)
        
        try:
            provider = NativePPTXProvider()
            extracted_doc = provider.extract_document(temp_path)
            
            # Check for chart/figure blocks
            figure_blocks = [b for b in extracted_doc.blocks if b.type == BlockType.FIGURE]
            assert len(figure_blocks) >= 1
            
            # Verify chart info extracted
            chart_block = figure_blocks[0]
            assert "Quarterly Sales" in chart_block.content
            assert "Chart Type:" in chart_block.content
            
        finally:
            temp_path.unlink()


if __name__ == "__main__":
    pytest.main([__file__, "-v"])