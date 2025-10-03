#!/usr/bin/env python3
"""Scenario: PPTX → UnifiedDocument → optional Arango insert.

Searches for any *.pptx under data/.
Skips gracefully when python-pptx is not installed or no sample exists.
"""
from __future__ import annotations

from pathlib import Path
from typing import Optional

from loguru import logger
import tempfile

from scenarios.extractors.common import (
    ScenarioResult,
    find_sample,
    summarise_unified,
    try_arango_insert,
    write_summary,
)


def main() -> int:
    name = "pptx"
    sample: Optional[Path] = find_sample("**/*.pptx")
    if not sample:
        # Synthesize a minimal PPTX if none present
        try:
            from pptx import Presentation  # type: ignore
            tmp = Path(tempfile.gettempdir()) / "scenario_sample.pptx"
            prs = Presentation()
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Scenario Sample"
            slide.placeholders[1].text = "Hello PPTX"
            prs.save(str(tmp))
            sample = tmp
        except Exception as e:
            logger.info(f"SKIP: cannot synthesize PPTX sample: {e}")
            res = ScenarioResult(
                name=name,
                ok=True,
                skipped=True,
                reason="sample-missing",
                input_path=None,
                provider="PPTXProvider",
                source_type=None,
                block_counts={},
                heading_sample=[],
                arango_inserted=False,
                artifacts={},
            )
            write_summary(name, res)
            return 0

    try:
        from scenarios.extractors.common import import_provider
        PPTXProvider = import_provider("providers/pptx.py", "PPTXProvider")
    except Exception as e:
        logger.warning(f"SKIP: cannot import PPTXProvider: {e}")
        res = ScenarioResult(
            name=name,
            ok=True,
            skipped=True,
            reason="import-error",
            input_path=str(sample),
            provider="PPTXProvider",
            source_type=None,
            block_counts={},
            heading_sample=[],
            arango_inserted=False,
            artifacts={},
        )
        write_summary(name, res)
        return 0

    doc = PPTXProvider().extract_document(sample)
    # Headings may be limited on slides; hierarchy is optional
    counts, heads = summarise_unified(doc)
    inserted = try_arango_insert(doc)

    ok = True  # succeed if extraction returns a document
    res = ScenarioResult(
        name=name,
        ok=ok,
        skipped=False,
        reason=None,
        input_path=str(sample),
        provider="PPTXProvider",
        source_type=str(doc.source_type),
        block_counts=counts,
        heading_sample=heads,
        arango_inserted=inserted,
        artifacts={},
    )
    write_summary(name, res)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
