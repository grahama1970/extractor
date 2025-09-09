#!/usr/bin/env python3
"""
PowerPoint Document Extractor Worker

Extracts content from PowerPoint presentations (PPTX) while preserving
slides, layouts, speaker notes, and embedded objects.
"""

import json
import asyncio
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple
from datetime import datetime
import hashlib
import base64

import typer
from loguru import logger
from rich.console import Console
from rich.table import Table
from rich.panel import Panel
from rich.progress import Progress, SpinnerColumn, TextColumn

try:
    from pptx import Presentation
    from pptx.slide import Slide
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available - PowerPoint support disabled")

# Configure logger
logger.remove()
logger.add(lambda msg: print(msg, end=""), level="INFO", format="{message}")

app = typer.Typer(help="Extract content from PowerPoint presentations")
console = Console()


class PowerPointExtractor:
    """Extract content from PowerPoint presentations."""
    
    def __init__(self):
        self.cache_dir = Path.home() / ".cache" / "extractor" / "pptx"
        self.cache_dir.mkdir(parents=True, exist_ok=True)
        
    async def extract_presentation(self,
                                 pptx_path: Path,
                                 extract_notes: bool = True,
                                 extract_images: bool = True,
                                 extract_charts: bool = True,
                                 max_image_size_mb: int = 10) -> Dict:
        """Extract content from PowerPoint presentation.
        
        Args:
            pptx_path: Path to PowerPoint file
            extract_notes: Whether to extract speaker notes
            extract_images: Whether to extract embedded images
            extract_charts: Whether to extract chart data
            max_image_size_mb: Maximum image size to embed
            
        Returns:
            Extracted content with metadata
        """
        # Validate path
        pptx_path = Path(pptx_path).resolve()
        if not pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
        if pptx_path.suffix.lower() not in ['.pptx', '.ppt', '.pptm']:
            raise ValueError(f"Not a PowerPoint file: {pptx_path}")
        
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PowerPoint extraction")
        
        # Check cache
        cache_key = self._get_cache_key(pptx_path)
        cached = await self._check_cache(cache_key)
        if cached:
            logger.info(f"Using cached extraction for {pptx_path.name}")
            return cached
        
        # Load presentation
        prs = Presentation(str(pptx_path))
        
        # Extract metadata
        metadata = self._extract_metadata(prs, pptx_path)
        
        # Extract slides
        slides = []
        for slide_idx, slide in enumerate(prs.slides):
            slide_data = await self._extract_slide(
                slide, slide_idx,
                extract_notes, extract_images, extract_charts,
                max_image_size_mb
            )
            slides.append(slide_data)
        
        result = {
            "slides": slides,
            "metadata": metadata,
            "statistics": self._calculate_statistics(slides)
        }
        
        # Cache result
        await self._cache_result(cache_key, result)
        
        return result
    
    async def _extract_slide(self,
                           slide: Slide,
                           slide_idx: int,
                           extract_notes: bool,
                           extract_images: bool,
                           extract_charts: bool,
                           max_image_size_mb: int) -> Dict:
        """Extract content from a single slide."""
        slide_data = {
            "number": slide_idx + 1,
            "title": self._get_slide_title(slide),
            "content": [],
            "shapes": [],
            "notes": "",
            "layout": slide.slide_layout.name if slide.slide_layout else "Custom",
            "shape_count": len(slide.shapes)
        }
        
        # Extract shapes
        for shape in slide.shapes:
            shape_data = await self._extract_shape(
                shape, extract_images, extract_charts, max_image_size_mb
            )
            if shape_data:
                slide_data["shapes"].append(shape_data)
        
        # Extract speaker notes
        if extract_notes and slide.has_notes_slide:
            notes_text = []
            for shape in slide.notes_slide.shapes:
                if shape.has_text_frame and shape.text:
                    notes_text.append(shape.text.strip())
            slide_data["notes"] = '\n'.join(notes_text)
        
        # Build slide content summary
        slide_data["content"] = self._build_slide_content(slide_data["shapes"])
        
        return slide_data
    
    def _get_slide_title(self, slide: Slide) -> str:
        """Extract slide title from title placeholder."""
        for shape in slide.shapes:
            if shape.is_placeholder:
                if shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                    if hasattr(shape, "text") and shape.text:
                        return shape.text.strip()
        return ""
    
    async def _extract_shape(self,
                           shape: Any,
                           extract_images: bool,
                           extract_charts: bool,
                           max_image_size_mb: int) -> Optional[Dict]:
        """Extract content from a shape."""
        shape_data = {
            "type": self._get_shape_type(shape),
            "position": {
                "left": shape.left if hasattr(shape, 'left') else None,
                "top": shape.top if hasattr(shape, 'top') else None,
                "width": shape.width if hasattr(shape, 'width') else None,
                "height": shape.height if hasattr(shape, 'height') else None
            }
        }
        
        # Handle different shape types
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Grouped shapes
            shape_data["group_items"] = []
            for sub_shape in shape.shapes:
                sub_data = await self._extract_shape(
                    sub_shape, extract_images, extract_charts, max_image_size_mb
                )
                if sub_data:
                    shape_data["group_items"].append(sub_data)
        
        elif shape.has_table:
            # Table
            shape_data["table"] = self._extract_table(shape.table)
        
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and extract_images:
            # Image
            shape_data["image"] = await self._extract_image(shape, max_image_size_mb)
        
        elif shape.has_chart and extract_charts:
            # Chart
            shape_data["chart"] = self._extract_chart(shape.chart)
        
        elif hasattr(shape, "text") and shape.text and shape.text.strip():
            # Text
            shape_data["text"] = shape.text.strip()
            if shape.has_text_frame:
                shape_data["paragraphs"] = self._extract_paragraphs(shape.text_frame)
        
        else:
            # Skip empty shapes
            return None
        
        return shape_data
    
    def _get_shape_type(self, shape: Any) -> str:
        """Determine shape type."""
        if shape.has_table:
            return "table"
        elif shape.has_chart:
            return "chart"
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return "group"
        elif hasattr(shape, "text") and shape.text:
            if shape.is_placeholder:
                placeholder_type = shape.placeholder_format.type
                if placeholder_type == PP_PLACEHOLDER.TITLE:
                    return "title"
                elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                    return "subtitle"
                elif placeholder_type == PP_PLACEHOLDER.BODY:
                    return "body"
            return "text"
        else:
            return shape.shape_type.name if hasattr(shape.shape_type, 'name') else "unknown"
    
    def _extract_table(self, table) -> Dict:
        """Extract table data."""
        table_data = {
            "rows": len(table.rows),
            "cols": len(table.columns),
            "cells": []
        }
        
        for row_idx, row in enumerate(table.rows):
            row_data = []
            for col_idx, cell in enumerate(row.cells):
                cell_text = cell.text.strip() if cell.text else ""
                row_data.append(cell_text)
            table_data["cells"].append(row_data)
        
        # Detect headers (first row)
        if table_data["cells"]:
            table_data["headers"] = table_data["cells"][0]
        
        return table_data
    
    async def _extract_image(self, shape: Any, max_size_mb: int) -> Dict:
        """Extract image data."""
        try:
            image = shape.image
            image_bytes = image.blob
            
            image_data = {
                "format": image.ext or 'unknown',
                "size_bytes": len(image_bytes),
                "filename": shape.name or f"image_{shape.shape_id}"
            }
            
            # Check size limit
            max_bytes = max_size_mb * 1024 * 1024
            if len(image_bytes) <= max_bytes:
                # Encode as base64
                image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                mime_type = f"image/{image.ext}" if image.ext else "image/png"
                image_data["data_uri"] = f"data:{mime_type};base64,{image_base64}"
            else:
                image_data["skipped"] = "Image too large"
            
            return image_data
            
        except Exception as e:
            logger.warning(f"Failed to extract image: {e}")
            return {"error": str(e)}
    
    def _extract_chart(self, chart) -> Dict:
        """Extract chart information."""
        chart_data = {
            "type": "unknown",
            "title": "",
            "series_count": 0
        }
        
        try:
            # Chart type
            if hasattr(chart, 'chart_type'):
                chart_data["type"] = str(chart.chart_type)
            
            # Title
            if hasattr(chart, 'has_title') and chart.has_title:
                if hasattr(chart, 'chart_title') and chart.chart_title.text_frame:
                    chart_data["title"] = chart.chart_title.text_frame.text
            
            # Series count
            if hasattr(chart, 'series'):
                chart_data["series_count"] = len(chart.series)
                
        except Exception as e:
            logger.debug(f"Chart extraction error: {e}")
        
        return chart_data
    
    def _extract_paragraphs(self, text_frame) -> List[Dict]:
        """Extract paragraph formatting."""
        paragraphs = []
        
        for para in text_frame.paragraphs:
            para_data = {
                "text": para.text,
                "level": para.level,
                "runs": []
            }
            
            # Extract runs with formatting
            for run in para.runs:
                run_data = {
                    "text": run.text,
                    "bold": run.font.bold if run.font else None,
                    "italic": run.font.italic if run.font else None,
                    "underline": run.font.underline if run.font else None
                }
                para_data["runs"].append(run_data)
            
            paragraphs.append(para_data)
        
        return paragraphs
    
    def _build_slide_content(self, shapes: List[Dict]) -> List[str]:
        """Build content summary from shapes."""
        content = []
        
        for shape in shapes:
            if shape["type"] == "title" and "text" in shape:
                content.append(f"Title: {shape['text']}")
            elif shape["type"] == "subtitle" and "text" in shape:
                content.append(f"Subtitle: {shape['text']}")
            elif shape["type"] == "body" and "text" in shape:
                content.append(shape["text"])
            elif shape["type"] == "text" and "text" in shape:
                content.append(shape["text"])
            elif shape["type"] == "table":
                table = shape.get("table", {})
                content.append(f"Table: {table.get('rows', 0)}x{table.get('cols', 0)}")
            elif shape["type"] == "chart":
                chart = shape.get("chart", {})
                title = chart.get("title", "Untitled")
                content.append(f"Chart: {title}")
            elif shape["type"] == "image":
                content.append("Image")
        
        return content
    
    def _extract_metadata(self, prs: Any, file_path: Path) -> Dict:
        """Extract presentation metadata."""
        metadata = {
            "file_name": file_path.name,
            "file_size": file_path.stat().st_size,
            "slide_count": len(prs.slides),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
            "aspect_ratio": f"{prs.slide_width}x{prs.slide_height}"
        }
        
        # Core properties
        if hasattr(prs, 'core_properties'):
            props = prs.core_properties
            metadata.update({
                "title": props.title or "",
                "author": props.author or "",
                "created": props.created.isoformat() if props.created else None,
                "modified": props.modified.isoformat() if props.modified else None,
                "subject": props.subject or "",
                "keywords": props.keywords or "",
                "comments": props.comments or ""
            })
        
        return metadata
    
    def _calculate_statistics(self, slides: List[Dict]) -> Dict:
        """Calculate presentation statistics."""
        stats = {
            "total_slides": len(slides),
            "slides_with_notes": 0,
            "total_shapes": 0,
            "shape_types": {},
            "total_tables": 0,
            "total_charts": 0,
            "total_images": 0,
            "total_text_blocks": 0
        }
        
        for slide in slides:
            if slide.get("notes"):
                stats["slides_with_notes"] += 1
            
            stats["total_shapes"] += slide.get("shape_count", 0)
            
            for shape in slide.get("shapes", []):
                shape_type = shape.get("type", "unknown")
                stats["shape_types"][shape_type] = stats["shape_types"].get(shape_type, 0) + 1
                
                if shape_type == "table":
                    stats["total_tables"] += 1
                elif shape_type == "chart":
                    stats["total_charts"] += 1
                elif shape_type == "image":
                    stats["total_images"] += 1
                elif shape_type in ["text", "title", "subtitle", "body"]:
                    stats["total_text_blocks"] += 1
        
        return stats
    
    def _get_cache_key(self, file_path: Path) -> str:
        """Generate cache key."""
        stat = file_path.stat()
        data = f"{file_path.absolute()}:{stat.st_size}:{stat.st_mtime}"
        return hashlib.sha256(data.encode()).hexdigest()
    
    async def _check_cache(self, cache_key: str) -> Optional[Dict]:
        """Check cache for existing extraction."""
        cache_file = self.cache_dir / f"{cache_key}.json"
        if cache_file.exists():
            try:
                with open(cache_file) as f:
                    return json.load(f)
            except:
                pass
        return None
    
    async def _cache_result(self, cache_key: str, result: Dict):
        """Cache extraction result."""
        cache_file = self.cache_dir / f"{cache_key}.json"
        try:
            with open(cache_file, 'w') as f:
                json.dump(result, f, indent=2)
        except Exception as e:
            logger.warning(f"Failed to cache: {e}")
    
    def display_slide_preview(self, slide_data: Dict):
        """Display a preview of slide content."""
        title = slide_data.get("title", f"Slide {slide_data['number']}")
        panel_content = []
        
        # Add content
        content = slide_data.get("content", [])
        if content:
            panel_content.append("[bold]Content:[/bold]")
            for item in content[:5]:  # Show first 5 items
                if len(item) > 100:
                    item = item[:97] + "..."
                panel_content.append(f"  {item}")
            if len(content) > 5:
                panel_content.append(f"  ... and {len(content) - 5} more items")
        
        # Add notes preview
        notes = slide_data.get("notes", "")
        if notes:
            panel_content.append("\n[bold]Speaker Notes:[/bold]")
            notes_preview = notes[:200] + "..." if len(notes) > 200 else notes
            panel_content.append(f"  {notes_preview}")
        
        # Shape summary
        shape_count = slide_data.get("shape_count", 0)
        panel_content.append(f"\n[dim]{shape_count} shapes on slide[/dim]")
        
        console.print(Panel(
            "\n".join(panel_content),
            title=f"[bold cyan]{title}[/bold cyan]",
            border_style="cyan"
        ))


# Initialize extractor
extractor = PowerPointExtractor()


@app.command()
def extract(
    pptx_file: Path = typer.Argument(..., help="PowerPoint file to extract"),
    output: Optional[Path] = typer.Option(None, "--output", "-o", help="Output JSON file"),
    no_notes: bool = typer.Option(False, "--no-notes", help="Don't extract speaker notes"),
    no_images: bool = typer.Option(False, "--no-images", help="Don't extract images"),
    no_charts: bool = typer.Option(False, "--no-charts", help="Don't extract charts"),
    preview: bool = typer.Option(False, "--preview", "-p", help="Show slide previews"),
    slides: Optional[str] = typer.Option(None, "--slides", "-s", help="Specific slide numbers (e.g. 1,3-5)")
):
    """Extract content from PowerPoint presentation."""
    async def run():
        try:
            # Extract
            result = await extractor.extract_presentation(
                pptx_file,
                extract_notes=not no_notes,
                extract_images=not no_images,
                extract_charts=not no_charts
            )
            
            # Filter slides if requested
            if slides:
                slide_nums = _parse_slide_range(slides)
                result["slides"] = [
                    s for s in result["slides"] 
                    if s["number"] in slide_nums
                ]
            
            # Show preview if requested
            if preview:
                for slide in result["slides"][:5]:  # Show first 5 slides
                    extractor.display_slide_preview(slide)
                    console.print()
                
                if len(result["slides"]) > 5:
                    console.print(f"[dim]... and {len(result['slides']) - 5} more slides[/dim]")
            
            # Save or display summary
            if output:
                with open(output, 'w') as f:
                    json.dump(result, f, indent=2)
                console.print(f"[green] Saved to {output}[/green]")
            else:
                # Display summary
                stats = result["statistics"]
                metadata = result["metadata"]
                
                table = Table(title="PowerPoint Extraction Summary")
                table.add_column("Metric", style="cyan")
                table.add_column("Value", style="green")
                
                table.add_row("File", metadata["file_name"])
                table.add_row("Total Slides", str(stats["total_slides"]))
                table.add_row("Slides with Notes", str(stats["slides_with_notes"]))
                table.add_row("Total Shapes", str(stats["total_shapes"]))
                table.add_row("Tables", str(stats["total_tables"]))
                table.add_row("Charts", str(stats["total_charts"]))
                table.add_row("Images", str(stats["total_images"]))
                table.add_row("Text Blocks", str(stats["total_text_blocks"]))
                
                console.print(table)
        
        except Exception as e:
            console.print(f"[red]Error: {e}[/red]")
            raise typer.Exit(1)
    
    asyncio.run(run())


@app.command()
def analyze(
    pptx_file: Path = typer.Argument(..., help="PowerPoint file to analyze"),
    verbose: bool = typer.Option(False, "--verbose", "-v", help="Show detailed analysis")
):
    """Analyze PowerPoint structure and content."""
    async def run():
        try:
            result = await extractor.extract_presentation(pptx_file)
            
            console.print(f"\n[bold]Analysis of {pptx_file.name}[/bold]\n")
            
            # Metadata
            metadata = result["metadata"]
            console.print(f"Title: {metadata.get('title', 'Untitled')}")
            console.print(f"Author: {metadata.get('author', 'Unknown')}")
            console.print(f"Slides: {metadata['slide_count']}")
            console.print(f"Aspect Ratio: {metadata['aspect_ratio']}")
            
            if metadata.get("created"):
                console.print(f"Created: {metadata['created']}")
            if metadata.get("modified"):
                console.print(f"Modified: {metadata['modified']}")
            
            # Statistics
            stats = result["statistics"]
            console.print(f"\n[bold]Content Statistics:[/bold]")
            console.print(f"Slides with notes: {stats['slides_with_notes']}/{stats['total_slides']}")
            console.print(f"Total shapes: {stats['total_shapes']}")
            
            if stats["shape_types"]:
                console.print("\n[bold]Shape Types:[/bold]")
                for shape_type, count in sorted(stats["shape_types"].items()):
                    console.print(f"  {shape_type}: {count}")
            
            if verbose:
                # Show slide details
                console.print("\n[bold]Slide Details:[/bold]")
                for slide in result["slides"]:
                    title = slide.get("title", f"Slide {slide['number']}")
                    console.print(f"\n{slide['number']}. {title}")
                    console.print(f"   Layout: {slide['layout']}")
                    console.print(f"   Shapes: {slide['shape_count']}")
                    if slide.get("notes"):
                        console.print(f"   Has speaker notes: Yes")
        
        except Exception as e:
            console.print(f"[red]Error: {e}[/red]")
            raise typer.Exit(1)
    
    asyncio.run(run())


def _parse_slide_range(slides_str: str) -> List[int]:
    """Parse slide range string (e.g. '1,3-5,7')."""
    slide_nums = []
    for part in slides_str.split(','):
        if '-' in part:
            start, end = part.split('-')
            slide_nums.extend(range(int(start), int(end) + 1))
        else:
            slide_nums.append(int(part))
    return slide_nums


# Worker functions
async def working_usage():
    """Demonstrate PowerPoint extraction."""
    logger.info("Testing PowerPoint extraction...")
    
    if not PPTX_AVAILABLE:
        logger.warning("python-pptx not installed - skipping demo")
        return
    
    # Create test presentation
    from pptx import Presentation
    from pptx.util import Inches
    
    prs = Presentation()
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "Test Presentation"
    subtitle.text = "Extraction Demo"
    
    # Content slide
    content_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = content_slide.shapes.title
    content = content_slide.placeholders[1]
    title.text = "Key Points"
    
    # Add bullet points
    tf = content.text_frame
    tf.text = "First point"
    p = tf.add_paragraph()
    p.text = "Second point"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Sub-point"
    p.level = 1
    
    # Add speaker notes
    notes_slide = content_slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = "Remember to emphasize the importance of these points."
    
    # Save test file
    test_file = Path("/tmp/test_presentation.pptx")
    prs.save(test_file)
    
    # Extract
    result = await extractor.extract_presentation(test_file)
    
    logger.info(f"\nExtraction complete:")
    logger.info(f"  Slides: {result['statistics']['total_slides']}")
    logger.info(f"  Total shapes: {result['statistics']['total_shapes']}")
    logger.info(f"  Slides with notes: {result['statistics']['slides_with_notes']}")
    
    # Show slide content
    for slide in result["slides"]:
        logger.info(f"\nSlide {slide['number']}: {slide['title']}")
        for content in slide["content"]:
            logger.info(f"  - {content}")
        if slide["notes"]:
            logger.info(f"  Notes: {slide['notes']}")


async def debug_function():
    """Test edge cases in PowerPoint extraction."""
    logger.info("Testing PowerPoint edge cases...")
    
    if not PPTX_AVAILABLE:
        logger.warning("python-pptx not installed - skipping tests")
        return
    
    from pptx import Presentation
    from pptx.util import Inches
    
    prs = Presentation()
    
    # Slide with table
    table_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    shapes = table_slide.shapes
    
    # Add table
    rows, cols = 3, 3
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(4)
    height = Inches(2)
    
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    # Fill table
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(0, 2).text = "Header 3"
    
    for row in range(1, rows):
        for col in range(cols):
            table.cell(row, col).text = f"R{row}C{col}"
    
    # Slide with grouped shapes
    group_slide = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = group_slide.shapes
    
    # Add some shapes (can't actually group in python-pptx, but we can test shape detection)
    shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text = "Text Box"
    shapes.add_connector(1, Inches(1), Inches(2.5), Inches(3), Inches(2.5))
    
    # Save
    test_file = Path("/tmp/test_complex.pptx")
    prs.save(test_file)
    
    # Extract
    result = await extractor.extract_presentation(test_file)
    
    logger.info(f"\nComplex presentation extracted:")
    logger.info(f"  Tables found: {result['statistics']['total_tables']}")
    logger.info(f"  Shape types: {result['statistics']['shape_types']}")
    
    # Check table extraction
    for slide in result["slides"]:
        for shape in slide["shapes"]:
            if shape["type"] == "table":
                table_data = shape.get("table", {})
                logger.info(f"  Found table: {table_data['rows']}x{table_data['cols']}")
                if table_data.get("headers"):
                    logger.info(f"  Headers: {table_data['headers']}")


if __name__ == "__main__":
    import sys
    
    if len(sys.argv) > 1 and sys.argv[1] == "working_usage":
        asyncio.run(working_usage())
    elif len(sys.argv) > 1 and sys.argv[1] == "debug":
        asyncio.run(debug_function())
    else:
        app()