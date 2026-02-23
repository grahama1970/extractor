#!/usr/bin/env python3
"""
generate_multiformat_fixture.py

Generates deterministic test fixtures in ALL supported formats from a single source spec.
This enables true cross-format validation since all formats contain identical content.

Supported outputs:
  - PDF (via PyMuPDF)
  - HTML
  - Markdown
  - XML
  - RST (reStructuredText)
  - DOCX (via python-docx)
  - PPTX (via python-pptx)
  - XLSX (via openpyxl)
  - EPUB (via ebooklib)
  - PNG (rendered image via PyMuPDF)
  - JSON (structured ground truth)

Usage:
    python generate_multiformat_fixture.py --config spec.yml --output-dir fixtures/my_test/
    python generate_multiformat_fixture.py --config spec.yml --output-dir fixtures/my_test/ --formats pdf,html,md
    python generate_multiformat_fixture.py --config spec.yml --output-dir fixtures/my_test/ --formats all
"""

import argparse
import json
import yaml
import fitz  # PyMuPDF
from pathlib import Path
from typing import Dict, List, Any
from dataclasses import dataclass, field
from datetime import datetime

# Optional imports for DOCX
try:
    from docx import Document as DocxDocument

    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# Optional imports for PPTX
try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Optional imports for XLSX
try:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

# Optional imports for EPUB
try:
    from ebooklib import epub

    HAS_EPUB = True
except ImportError:
    HAS_EPUB = False


@dataclass
class ContentBlock:
    """Represents a single content block that can be rendered to any format."""

    block_type: str  # "heading", "paragraph", "table", "list", "requirement", "figure_placeholder"
    level: int = 1  # For headings: 1-6
    text: str = ""
    items: List[str] = field(default_factory=list)  # For lists
    rows: List[Dict[str, str]] = field(default_factory=list)  # For tables
    columns: List[str] = field(default_factory=list)  # Table column headers
    req_id: str = ""  # For requirements
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class DocumentSpec:
    """Specification for a multi-format document."""

    title: str
    blocks: List[ContentBlock]
    metadata: Dict[str, Any] = field(default_factory=dict)


def parse_spec(config_path: Path) -> DocumentSpec:
    """Parse a YAML config into a DocumentSpec."""
    config = yaml.safe_load(config_path.read_text())

    blocks = []
    for item in config.get("content", []):
        block_type = item.get("type", "paragraph")

        if block_type == "heading":
            blocks.append(
                ContentBlock(
                    block_type="heading", level=item.get("level", 1), text=item.get("text", "")
                )
            )
        elif block_type == "paragraph":
            blocks.append(ContentBlock(block_type="paragraph", text=item.get("text", "")))
        elif block_type == "requirement":
            blocks.append(
                ContentBlock(
                    block_type="requirement",
                    req_id=item.get("id", ""),
                    text=item.get("text", ""),
                    metadata={"strict": item.get("strict", True)},
                )
            )
        elif block_type == "list":
            blocks.append(ContentBlock(block_type="list", items=item.get("items", [])))
        elif block_type == "table":
            blocks.append(
                ContentBlock(
                    block_type="table",
                    columns=item.get("columns", []),
                    rows=item.get("rows", []),
                    text=item.get("caption", ""),
                )
            )
        elif block_type == "figure":
            blocks.append(
                ContentBlock(
                    block_type="figure_placeholder", text=item.get("caption", "Figure placeholder")
                )
            )

    return DocumentSpec(
        title=config.get("title", "Untitled"), blocks=blocks, metadata=config.get("metadata", {})
    )


class MultiFormatGenerator:
    """Generates documents in multiple formats from a single spec."""

    def __init__(self, spec: DocumentSpec):
        self.spec = spec
        self.ground_truth: List[Dict[str, Any]] = []

    def generate_all(self, output_dir: Path, formats: List[str] = None, name: str = "fixture"):
        """Generate all requested formats."""
        output_dir.mkdir(parents=True, exist_ok=True)

        all_formats = [
            "pdf",
            "html",
            "md",
            "xml",
            "rst",
            "docx",
            "pptx",
            "xlsx",
            "epub",
            "png",
            "json",
        ]
        formats = formats or all_formats

        results = {}
        skipped = []

        for fmt in formats:
            if fmt == "pdf":
                results["pdf"] = self._generate_pdf(output_dir / f"{name}.pdf")
            elif fmt == "html":
                results["html"] = self._generate_html(output_dir / f"{name}.html")
            elif fmt == "md":
                results["md"] = self._generate_markdown(output_dir / f"{name}.md")
            elif fmt == "xml":
                results["xml"] = self._generate_xml(output_dir / f"{name}.xml")
            elif fmt == "rst":
                results["rst"] = self._generate_rst(output_dir / f"{name}.rst")
            elif fmt == "docx":
                if HAS_DOCX:
                    results["docx"] = self._generate_docx(output_dir / f"{name}.docx")
                else:
                    skipped.append("docx (python-docx not installed)")
            elif fmt == "pptx":
                if HAS_PPTX:
                    results["pptx"] = self._generate_pptx(output_dir / f"{name}.pptx")
                else:
                    skipped.append("pptx (python-pptx not installed)")
            elif fmt == "xlsx":
                if HAS_XLSX:
                    results["xlsx"] = self._generate_xlsx(output_dir / f"{name}.xlsx")
                else:
                    skipped.append("xlsx (openpyxl not installed)")
            elif fmt == "epub":
                if HAS_EPUB:
                    results["epub"] = self._generate_epub(output_dir / f"{name}.epub")
                else:
                    skipped.append("epub (ebooklib not installed)")
            elif fmt == "png":
                # PNG is generated from PDF, so ensure PDF exists
                pdf_path = output_dir / f"{name}.pdf"
                if not pdf_path.exists():
                    self._generate_pdf(pdf_path)
                results["png"] = self._generate_png(pdf_path, output_dir / f"{name}.png")
            elif fmt == "json":
                results["json"] = self._generate_json(output_dir / f"{name}_expected.json")

        # Always generate ground truth JSON
        if "json" not in formats:
            self._generate_json(output_dir / f"{name}_expected.json")

        if skipped:
            print(f"Skipped formats: {', '.join(skipped)}")

        return results

    def _generate_pdf(self, path: Path) -> Path:
        """Generate PDF using PyMuPDF."""
        doc = fitz.open()
        page = doc.new_page()
        y = 50

        # Title
        page.insert_text((50, y), self.spec.title, fontsize=18, fontname="Helvetica-Bold")
        y += 40

        for block in self.spec.blocks:
            # Check if we need a new page
            if y > 750:
                page = doc.new_page()
                y = 50

            if block.block_type == "heading":
                fontsize = {1: 16, 2: 14, 3: 12, 4: 11, 5: 10, 6: 10}.get(block.level, 12)
                page.insert_text((50, y), block.text, fontsize=fontsize, fontname="Helvetica-Bold")
                y += fontsize + 10

            elif block.block_type == "paragraph":
                # Use textbox for wrapping
                rect = fitz.Rect(50, y, 545, y + 100)
                rc = page.insert_textbox(rect, block.text, fontsize=11)
                y += max(20, -rc + 15)

            elif block.block_type == "requirement":
                req_text = f"{block.req_id}: {block.text}"
                page.insert_text((50, y), req_text, fontsize=11)
                y += 20
                self.ground_truth.append(
                    {
                        "id": block.req_id,
                        "type": "Requirement",
                        "text": block.text,
                        "strict_verification": block.metadata.get("strict", True),
                    }
                )

            elif block.block_type == "list":
                for item in block.items:
                    page.insert_text((60, y), f"• {item}", fontsize=11)
                    y += 18

            elif block.block_type == "table":
                # Table caption
                if block.text:
                    page.insert_text((50, y), block.text, fontsize=10, fontname="Helvetica-Oblique")
                    y += 18

                # Draw proper table grid for Camelot lattice detection
                num_cols = len(block.columns)
                num_rows = len(block.rows) + 1  # +1 for header
                col_width = 480 // max(num_cols, 1)
                row_height = 20
                table_left = 50
                table_right = table_left + col_width * num_cols
                table_top = y
                table_bottom = y + row_height * num_rows

                # Draw outer border (thick line for better detection)
                page.draw_rect(
                    fitz.Rect(table_left, table_top, table_right, table_bottom),
                    color=(0, 0, 0),
                    width=1.5,
                )

                # Draw horizontal lines for all rows
                for row_idx in range(num_rows + 1):
                    row_y = table_top + row_idx * row_height
                    page.draw_line(
                        (table_left, row_y), (table_right, row_y), color=(0, 0, 0), width=1
                    )

                # Draw vertical lines for all columns
                for col_idx in range(num_cols + 1):
                    col_x = table_left + col_idx * col_width
                    page.draw_line(
                        (col_x, table_top), (col_x, table_bottom), color=(0, 0, 0), width=1
                    )

                # Header row background
                page.draw_rect(
                    fitz.Rect(table_left, table_top, table_right, table_top + row_height),
                    color=(0.85, 0.85, 0.85),
                    fill=(0.85, 0.85, 0.85),
                )
                # Re-draw header borders on top of fill
                page.draw_line(
                    (table_left, table_top), (table_right, table_top), color=(0, 0, 0), width=1
                )
                page.draw_line(
                    (table_left, table_top + row_height),
                    (table_right, table_top + row_height),
                    color=(0, 0, 0),
                    width=1,
                )

                # Header text
                for i, col in enumerate(block.columns):
                    page.insert_text(
                        (table_left + 5 + i * col_width, table_top + 14),
                        col,
                        fontsize=10,
                        fontname="Helvetica-Bold",
                    )

                # Data rows
                for row_idx, row in enumerate(block.rows):
                    row_y = table_top + (row_idx + 1) * row_height
                    for i, col in enumerate(block.columns):
                        page.insert_text(
                            (table_left + 5 + i * col_width, row_y + 14),
                            str(row.get(col, "")),
                            fontsize=10,
                        )

                y = table_bottom + 10

                self.ground_truth.append(
                    {
                        "id": f"TBL-{len([g for g in self.ground_truth if g['type'] == 'Table']) + 1:03d}",
                        "type": "Table",
                        "caption": block.text,
                        "columns": block.columns,
                        "rows": block.rows,
                    }
                )

            elif block.block_type == "figure_placeholder":
                page.draw_rect(fitz.Rect(50, y, 300, y + 100), color=(0.7, 0.7, 0.7))
                page.insert_text(
                    (60, y + 50), f"[{block.text}]", fontsize=10, fontname="Helvetica-Oblique"
                )
                y += 120

        doc.save(path)
        print(f"Generated PDF: {path}")
        return path

    def _generate_html(self, path: Path) -> Path:
        """Generate HTML document."""
        lines = [
            "<!DOCTYPE html>",
            "<html>",
            "<head>",
            f"    <title>{self.spec.title}</title>",
            '    <meta charset="UTF-8">',
            "</head>",
            "<body>",
            f"<h1>{self.spec.title}</h1>",
        ]

        for block in self.spec.blocks:
            if block.block_type == "heading":
                lines.append(f"<h{block.level}>{block.text}</h{block.level}>")

            elif block.block_type == "paragraph":
                lines.append(f"<p>{block.text}</p>")

            elif block.block_type == "requirement":
                lines.append(f"<p><strong>{block.req_id}</strong>: {block.text}</p>")

            elif block.block_type == "list":
                lines.append("<ul>")
                for item in block.items:
                    lines.append(f"    <li>{item}</li>")
                lines.append("</ul>")

            elif block.block_type == "table":
                if block.text:
                    lines.append(f"<p><em>{block.text}</em></p>")
                lines.append('<table border="1">')
                lines.append("    <thead><tr>")
                for col in block.columns:
                    lines.append(f"        <th>{col}</th>")
                lines.append("    </tr></thead>")
                lines.append("    <tbody>")
                for row in block.rows:
                    lines.append("        <tr>")
                    for col in block.columns:
                        lines.append(f'            <td>{row.get(col, "")}</td>')
                    lines.append("        </tr>")
                lines.append("    </tbody>")
                lines.append("</table>")

            elif block.block_type == "figure_placeholder":
                lines.append(f"<p><em>[{block.text}]</em></p>")

        lines.extend(["</body>", "</html>"])
        path.write_text("\n".join(lines))
        print(f"Generated HTML: {path}")
        return path

    def _generate_markdown(self, path: Path) -> Path:
        """Generate Markdown document."""
        lines = [f"# {self.spec.title}", ""]

        for block in self.spec.blocks:
            if block.block_type == "heading":
                prefix = "#" * (block.level + 1)  # +1 because title is H1
                lines.append(f"{prefix} {block.text}")
                lines.append("")

            elif block.block_type == "paragraph":
                lines.append(block.text)
                lines.append("")

            elif block.block_type == "requirement":
                lines.append(f"**{block.req_id}**: {block.text}")
                lines.append("")

            elif block.block_type == "list":
                for item in block.items:
                    lines.append(f"- {item}")
                lines.append("")

            elif block.block_type == "table":
                if block.text:
                    lines.append(f"*{block.text}*")
                    lines.append("")
                # Header
                lines.append("| " + " | ".join(block.columns) + " |")
                lines.append("|" + "|".join(["---"] * len(block.columns)) + "|")
                # Rows
                for row in block.rows:
                    cells = [str(row.get(col, "")) for col in block.columns]
                    lines.append("| " + " | ".join(cells) + " |")
                lines.append("")

            elif block.block_type == "figure_placeholder":
                lines.append(f"*[{block.text}]*")
                lines.append("")

        path.write_text("\n".join(lines))
        print(f"Generated Markdown: {path}")
        return path

    def _generate_xml(self, path: Path) -> Path:
        """Generate XML document."""
        lines = [
            '<?xml version="1.0" encoding="UTF-8"?>',
            "<document>",
            f"    <title>{self.spec.title}</title>",
            "    <metadata>",
            f"        <generated>{datetime.now().isoformat()}</generated>",
            "    </metadata>",
            "    <content>",
        ]

        for block in self.spec.blocks:
            if block.block_type == "heading":
                lines.append(f'        <section level="{block.level}" title="{block.text}">')
                lines.append("        </section>")

            elif block.block_type == "paragraph":
                lines.append(f"        <paragraph>{block.text}</paragraph>")

            elif block.block_type == "requirement":
                lines.append(f'        <requirement id="{block.req_id}">{block.text}</requirement>')

            elif block.block_type == "list":
                lines.append("        <list>")
                for item in block.items:
                    lines.append(f"            <item>{item}</item>")
                lines.append("        </list>")

            elif block.block_type == "table":
                lines.append(f'        <table caption="{block.text}">')
                lines.append("            <header>")
                for col in block.columns:
                    lines.append(f"                <column>{col}</column>")
                lines.append("            </header>")
                for row in block.rows:
                    lines.append("            <row>")
                    for col in block.columns:
                        lines.append(f'                <cell>{row.get(col, "")}</cell>')
                    lines.append("            </row>")
                lines.append("        </table>")

            elif block.block_type == "figure_placeholder":
                lines.append(f'        <figure caption="{block.text}"/>')

        lines.extend(["    </content>", "</document>"])
        path.write_text("\n".join(lines))
        print(f"Generated XML: {path}")
        return path

    def _generate_rst(self, path: Path) -> Path:
        """Generate reStructuredText document."""
        lines = []

        # Title with overline
        title_line = "=" * len(self.spec.title)
        lines.extend([title_line, self.spec.title, title_line, ""])

        for block in self.spec.blocks:
            if block.block_type == "heading":
                underline_chars = {1: "=", 2: "-", 3: "~", 4: "^", 5: '"', 6: "'"}
                char = underline_chars.get(block.level, "-")
                lines.append(block.text)
                lines.append(char * len(block.text))
                lines.append("")

            elif block.block_type == "paragraph":
                lines.append(block.text)
                lines.append("")

            elif block.block_type == "requirement":
                lines.append(f"**{block.req_id}**: {block.text}")
                lines.append("")

            elif block.block_type == "list":
                for item in block.items:
                    lines.append(f"* {item}")
                lines.append("")

            elif block.block_type == "table":
                if block.text:
                    lines.append(f"*{block.text}*")
                    lines.append("")

                # RST list-table directive
                lines.append(".. list-table::")
                lines.append("   :header-rows: 1")
                lines.append("")

                # Header row
                lines.append("   * - " + "\n     - ".join(block.columns))

                # Data rows
                for row in block.rows:
                    cells = [str(row.get(col, "")) for col in block.columns]
                    lines.append("   * - " + "\n     - ".join(cells))

                lines.append("")

            elif block.block_type == "figure_placeholder":
                lines.append(f"*[{block.text}]*")
                lines.append("")

        path.write_text("\n".join(lines))
        print(f"Generated RST: {path}")
        return path

    def _generate_docx(self, path: Path) -> Path:
        """Generate DOCX document."""
        if not HAS_DOCX:
            return None

        doc = DocxDocument()

        # Title
        doc.add_heading(self.spec.title, level=0)

        for block in self.spec.blocks:
            if block.block_type == "heading":
                doc.add_heading(block.text, level=block.level)

            elif block.block_type == "paragraph":
                doc.add_paragraph(block.text)

            elif block.block_type == "requirement":
                p = doc.add_paragraph()
                p.add_run(f"{block.req_id}: ").bold = True
                p.add_run(block.text)

            elif block.block_type == "list":
                for item in block.items:
                    doc.add_paragraph(item, style="List Bullet")

            elif block.block_type == "table":
                if block.text:
                    p = doc.add_paragraph()
                    p.add_run(block.text).italic = True

                # Create table
                table = doc.add_table(rows=1, cols=len(block.columns))
                table.style = "Table Grid"

                # Header
                header_cells = table.rows[0].cells
                for i, col in enumerate(block.columns):
                    header_cells[i].text = col

                # Data rows
                for row_data in block.rows:
                    row = table.add_row().cells
                    for i, col in enumerate(block.columns):
                        row[i].text = str(row_data.get(col, ""))

                doc.add_paragraph()  # Spacing

            elif block.block_type == "figure_placeholder":
                p = doc.add_paragraph()
                p.add_run(f"[{block.text}]").italic = True

        doc.save(path)
        print(f"Generated DOCX: {path}")
        return path

    def _generate_pptx(self, path: Path) -> Path:
        """Generate PowerPoint presentation."""
        if not HAS_PPTX:
            return None

        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = self.spec.title
        subtitle.text = f"Generated: {datetime.now().strftime('%Y-%m-%d')}"

        # Content slides
        content_slide_layout = prs.slide_layouts[1]  # Title and Content

        current_slide = None
        current_content = []

        for block in self.spec.blocks:
            if block.block_type == "heading" and block.level <= 2:
                # New slide for major headings
                if current_slide is not None and current_content:
                    # Add accumulated content to previous slide
                    body = current_slide.shapes.placeholders[1]
                    tf = body.text_frame
                    for i, text in enumerate(current_content):
                        if i == 0:
                            tf.text = text
                        else:
                            p = tf.add_paragraph()
                            p.text = text

                current_slide = prs.slides.add_slide(content_slide_layout)
                current_slide.shapes.title.text = block.text
                current_content = []

            elif block.block_type == "paragraph":
                current_content.append(block.text)

            elif block.block_type == "requirement":
                current_content.append(f"{block.req_id}: {block.text}")

            elif block.block_type == "list":
                for item in block.items:
                    current_content.append(f"  - {item}")

            elif block.block_type == "table":
                # Create actual table shape on current slide
                if current_slide is not None:
                    # First add any pending content
                    if current_content:
                        body = current_slide.shapes.placeholders[1]
                        tf = body.text_frame
                        for i, text in enumerate(current_content[:8]):
                            if i == 0:
                                tf.text = text
                            else:
                                p = tf.add_paragraph()
                                p.text = text
                        current_content = []

                    # Add table shape
                    from pptx.util import Inches as PptxInches

                    rows = len(block.rows) + 1  # +1 for header
                    cols = len(block.columns)
                    left = PptxInches(0.5)
                    top = PptxInches(3.5)
                    width = PptxInches(9)
                    height = PptxInches(0.5 * rows)

                    table_shape = current_slide.shapes.add_table(
                        rows, cols, left, top, width, height
                    )
                    tbl = table_shape.table

                    # Header row
                    for col_idx, col_name in enumerate(block.columns):
                        tbl.cell(0, col_idx).text = col_name

                    # Data rows
                    for row_idx, row_data in enumerate(block.rows):
                        for col_idx, col_name in enumerate(block.columns):
                            tbl.cell(row_idx + 1, col_idx).text = str(row_data.get(col_name, ""))

        # Finalize last slide
        if current_slide is not None and current_content:
            body = current_slide.shapes.placeholders[1]
            tf = body.text_frame
            for i, text in enumerate(current_content[:10]):  # Limit lines per slide
                if i == 0:
                    tf.text = text
                else:
                    p = tf.add_paragraph()
                    p.text = text

        prs.save(path)
        print(f"Generated PPTX: {path}")
        return path

    def _generate_xlsx(self, path: Path) -> Path:
        """Generate Excel spreadsheet with document content."""
        if not HAS_XLSX:
            return None

        wb = Workbook()

        # Document overview sheet
        ws = wb.active
        ws.title = "Overview"
        ws["A1"] = "Document Title"
        ws["B1"] = self.spec.title
        ws["A1"].font = Font(bold=True)

        row = 3
        ws[f"A{row}"] = "Content Summary"
        ws[f"A{row}"].font = Font(bold=True)
        row += 1

        for block in self.spec.blocks:
            if block.block_type == "heading":
                ws[f"A{row}"] = f"{'  ' * (block.level - 1)}{block.text}"
                ws[f"A{row}"].font = Font(bold=True)
            elif block.block_type == "paragraph":
                ws[f"A{row}"] = block.text[:100] + ("..." if len(block.text) > 100 else "")
            elif block.block_type == "requirement":
                ws[f"A{row}"] = block.req_id
                ws[f"B{row}"] = block.text
                ws[f"A{row}"].font = Font(bold=True)
            elif block.block_type == "list":
                for item in block.items:
                    ws[f"A{row}"] = f"  - {item}"
                    row += 1
                row -= 1  # Adjust for the extra increment
            row += 1

        # Create separate sheets for each table
        table_num = 0
        for block in self.spec.blocks:
            if block.block_type == "table":
                table_num += 1
                ws_table = wb.create_sheet(title=f"Table_{table_num}")

                # Caption
                if block.text:
                    ws_table["A1"] = block.text
                    ws_table["A1"].font = Font(italic=True)
                    start_row = 3
                else:
                    start_row = 1

                # Header
                for col_idx, col_name in enumerate(block.columns, 1):
                    cell = ws_table.cell(row=start_row, column=col_idx, value=col_name)
                    cell.font = Font(bold=True)

                # Data rows
                for row_idx, row_data in enumerate(block.rows, start_row + 1):
                    for col_idx, col_name in enumerate(block.columns, 1):
                        ws_table.cell(row=row_idx, column=col_idx, value=row_data.get(col_name, ""))

        wb.save(path)
        print(f"Generated XLSX: {path}")
        return path

    def _generate_epub(self, path: Path) -> Path:
        """Generate EPUB ebook."""
        if not HAS_EPUB:
            return None

        book = epub.EpubBook()

        # Metadata
        book.set_identifier(f'id-{datetime.now().strftime("%Y%m%d%H%M%S")}')
        book.set_title(self.spec.title)
        book.set_language("en")
        book.add_author("Fixture Generator")

        # Build HTML content
        html_content = [f"<h1>{self.spec.title}</h1>"]

        for block in self.spec.blocks:
            if block.block_type == "heading":
                html_content.append(f"<h{block.level + 1}>{block.text}</h{block.level + 1}>")
            elif block.block_type == "paragraph":
                html_content.append(f"<p>{block.text}</p>")
            elif block.block_type == "requirement":
                html_content.append(f"<p><strong>{block.req_id}</strong>: {block.text}</p>")
            elif block.block_type == "list":
                html_content.append("<ul>")
                for item in block.items:
                    html_content.append(f"<li>{item}</li>")
                html_content.append("</ul>")
            elif block.block_type == "table":
                if block.text:
                    html_content.append(f"<p><em>{block.text}</em></p>")
                html_content.append('<table border="1">')
                html_content.append(
                    "<tr>" + "".join(f"<th>{c}</th>" for c in block.columns) + "</tr>"
                )
                for row in block.rows:
                    html_content.append(
                        "<tr>"
                        + "".join(f'<td>{row.get(c, "")}</td>' for c in block.columns)
                        + "</tr>"
                    )
                html_content.append("</table>")

        # Create chapter
        chapter = epub.EpubHtml(title=self.spec.title, file_name="content.xhtml", lang="en")
        chapter.content = "\n".join(html_content)
        book.add_item(chapter)

        # Add navigation
        book.toc = [epub.Link("content.xhtml", self.spec.title, "content")]
        book.add_item(epub.EpubNcx())
        book.add_item(epub.EpubNav())
        book.spine = ["nav", chapter]

        epub.write_epub(path, book, {})
        print(f"Generated EPUB: {path}")
        return path

    def _generate_png(self, pdf_path: Path, png_path: Path) -> Path:
        """Generate PNG image from first page of PDF."""
        doc = fitz.open(pdf_path)
        if len(doc) > 0:
            page = doc[0]
            # Render at 2x resolution for clarity
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat)
            pix.save(png_path)
            print(f"Generated PNG: {png_path}")
        doc.close()
        return png_path

    def _generate_json(self, path: Path) -> Path:
        """Generate ground truth JSON."""
        output = {
            "metadata": {
                "title": self.spec.title,
                "generated": datetime.now().isoformat(),
                "formats": [
                    "pdf",
                    "html",
                    "md",
                    "xml",
                    "rst",
                    "docx",
                    "pptx",
                    "xlsx",
                    "epub",
                    "png",
                ],
            },
            "metrics": {
                "requirement_count": len(
                    [g for g in self.ground_truth if g["type"] == "Requirement"]
                ),
                "table_count": len([g for g in self.ground_truth if g["type"] == "Table"]),
                "total_items": len(self.ground_truth),
            },
            "ground_truth": self.ground_truth,
        }

        path.write_text(json.dumps(output, indent=2))
        print(f"Generated JSON: {path}")
        return path


def main():
    parser = argparse.ArgumentParser(description="Generate multi-format test fixtures")
    parser.add_argument("--config", type=Path, required=True, help="YAML spec file")
    parser.add_argument("--output-dir", type=Path, required=True, help="Output directory")
    parser.add_argument("--name", type=str, default="fixture", help="Base filename")
    parser.add_argument(
        "--formats",
        type=str,
        default="all",
        help="Comma-separated formats: pdf,html,md,xml,rst,docx,pptx,xlsx,epub,png,json or 'all'",
    )

    args = parser.parse_args()

    if not args.config.exists():
        print(f"Config not found: {args.config}")
        return 1

    spec = parse_spec(args.config)
    generator = MultiFormatGenerator(spec)

    formats = None if args.formats == "all" else args.formats.split(",")
    generator.generate_all(args.output_dir, formats=formats, name=args.name)

    print(f"\nAll fixtures generated in: {args.output_dir}")
    return 0


if __name__ == "__main__":
    exit(main())
